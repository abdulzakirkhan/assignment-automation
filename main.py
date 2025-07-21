# libraries
import os
import tempfile
import uuid
import logging
from dotenv import load_dotenv
from fastapi import FastAPI, Form, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware 
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import requests
import openai
from google.cloud import vision
import json
from docx import Document
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io
import shutil
import requests
from fastapi import Request
from io import BytesIO
import docx2txt
import os
import subprocess
import aspose.words as aw
from pdf2image import convert_from_path
from fastapi import HTTPException, Form, File, UploadFile
import zipfile
import tarfile
import rarfile
import csv
from datetime import datetime
from database import Database
import mysql.connector
from langchain.memory import ConversationBufferMemory
from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain_openai import ChatOpenAI
from openai import OpenAI
import re
import time
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import threading

load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure OpenAI API key securely
openai.api_key = os.getenv("OPENAI_API_KEY")
if not openai.api_key:
    logger.error("OpenAI API key not found. Please set OPENAI_API_KEY in environment variables.")
    raise EnvironmentError("OpenAI API key not found.")
openai_client = OpenAI(api_key=openai.api_key)

# Get the credentials path from the environment variable
google_credentials_path = os.getenv('GOOGLE_APPLICATION_CREDENTIALS')
logger.info(f"Service account path is: {google_credentials_path}")

# Check if the file exists
if google_credentials_path and os.path.exists(google_credentials_path):
    logger.info("service_account.json file found on disk!")
    os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = google_credentials_path
else:
    logger.warning("service_account.json file not found. OCR may not work properly.")

app = FastAPI()
client = vision.ImageAnnotatorClient()
db = Database() 

# All Allowed files
ALLOWED_FILE_TYPES = ["ppt", "pptx", "pdf", "doc", "docx", "jpeg", "jfif", "png", "zip", "jpg","rar", "tar", "tar.gz", "tar.bz2"]

# New Pydantic models for assignment generation
class AssignmentGenerationRequest(BaseModel):
    assignment_id: int

# ENHANCED SUBSECTION-WISE ASSIGNMENT GENERATOR FUNCTIONS

def extract_assignment_structure_advanced(assignment_text):
    """Advanced assignment structure extraction with universal compatibility"""
    
    system_prompt = """You are a universal assignment structure analyzer. Extract assignment structure from ANY academic document.

    Convert to this JSON format:
    {
      "assignment_title": "auto-detected title",
      "assignment_type": "auto-detected type (essay, report, briefing, etc.)",
      "total_word_count": "if mentioned",
      "assignment_sections": [
        {
          "section_id": 1,
          "section_title": "exact title from text",
          "word_count": "exact count if mentioned",
          "section_type": "auto-detected (introduction, main, conclusion, etc.)",
          "subsections": [
            {
              "subsection_id": "1.1", 
              "subsection_title": "exact subtitle",
              "content_requirements": "exact description/requirements",
              "word_count": "if specified",
              "theories_models": ["any theories mentioned"],
              "readings": ["any citations/readings mentioned"],
              "learning_outcomes": ["any learning outcomes mentioned"],
              "rubric_criteria": ["any rubric criteria mentioned"],
              "specific_instructions": ["any specific instructions"]
            }
          ]
        }
      ]
    }

    DYNAMIC EXTRACTION RULES:
    - ONLY create a 'subsections' array for a section if the original document or outline explicitly contains subsections, subheadings, or sub-parts for that section.
    - If a section does NOT have any explicit subsections in the outline or document, do NOT invent or create any subsections for it; either omit the 'subsections' key or set it as an empty array.
    - Do NOT split a section into subsections unless the source text clearly does so (e.g., with numbered subheadings, bullet points, or explicit sub-part titles).
    - Automatically detect assignment type and structure
    - Extract ALL sections regardless of format
    - Identify word counts, theories, readings anywhere in text
    - Preserve exact wording and requirements
    - Handle ANY academic assignment format
    - Include ALL details needed for content generation
    
    Extract structure from ANY assignment outline, syllabus, or academic document."""

    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Extract complete assignment structure from this text:\n\n{assignment_text}"}
            ],
            temperature=0.0
        )
        
        json_text = response.choices[0].message.content.strip()
        json_text = re.sub(r'^```json\s*\n?', '', json_text)
        json_text = re.sub(r'\n?\s*```$', '', json_text)
        
        start = json_text.find('{')
        end = json_text.rfind('}') + 1
        if start >= 0 and end > start:
            json_text = json_text[start:end]
        
        return json_text
        
    except Exception as e:
        logger.error(f"API Error: {e}")
        return None

def save_assignment_structure_advanced(json_text, assignment_id, output_dir="generated_assignments"):
    """Save assignment structure with advanced validation and display"""
    
    if not json_text:
        logger.error("No structure data received")
        return None
    
    try:
        data = json.loads(json_text)
        
        # Dynamic validation
        required_keys = ["assignment_sections"]
        for key in required_keys:
            if key not in data:
                logger.error(f"Missing required key: {key}")
                return None
        
        # Create assignment directory
        assignment_dir = os.path.join(output_dir, f"assignment_{assignment_id}")
        os.makedirs(assignment_dir, exist_ok=True)
        
        # Extract assignment details
        assignment_title = data.get("assignment_title", "Assignment")
        assignment_type = data.get("assignment_type", "Academic Assignment")
        total_word_count = data.get("total_word_count", "Not specified")
        sections = data["assignment_sections"]
        
        logger.info(f"✅ ASSIGNMENT DETECTED:")
        logger.info(f"   Title: {assignment_title}")
        logger.info(f"   Type: {assignment_type}")
        logger.info(f"   Total Words: {total_word_count}")
        logger.info(f"   Sections: {len(sections)}")
        
        # Show sections that will be generated
        for i, section in enumerate(sections, 1):
            title = section.get("section_title", f"Section {i}")
            word_count = section.get("word_count", "No limit")
            section_type = section.get("section_type", "unknown")
            subsections = section.get("subsections", [])
            
            logger.info(f"\n📋 Section {i}: {title}")
            logger.info(f"   Type: {section_type}")
            logger.info(f"   Word Count: {word_count}")
            logger.info(f"   Subsections: {len(subsections)}")
        
        # Save structure
        filename = os.path.join(assignment_dir, "assignment_structure.json")
        with open(filename, "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        
        logger.info(f"✅ Structure saved to '{filename}'")
        logger.info(f"🎯 Ready for subsection-wise content generation!")
        
        return data
        
    except json.JSONDecodeError as e:
        logger.error(f"JSON Error: {e}")
        return None

def determine_if_needs_citations(section_title, section_type):
    """Determine if a section should include citations based on title and type"""
    
    no_citation_keywords = [
        'summary', 'conclusion', 'concluding',
        'executive summary', 'abstract', 'preface', 'foreword',
        'acknowledgment', 'acknowledgement', 'table of contents',
        'contents', 'overview'
    ]
    
    intro_keywords = ['introduction', 'intro', 'background']
    title_lower = section_title.lower().strip()
    
    if any(keyword in title_lower for keyword in intro_keywords):
        return True
    
    if any(keyword in title_lower for keyword in no_citation_keywords):
        return False
    
    type_lower = section_type.lower().strip()
    if any(keyword in type_lower for keyword in no_citation_keywords):
        return False
        
    return True

def determine_if_needs_citations_subsection(section_title, subsection_title, section_type):
    """Determine if a subsection should include citations"""
    
    no_citation_keywords = [
        'summary', 'conclusion', 'concluding',
        'executive summary', 'abstract', 'preface', 'foreword',
        'acknowledgment', 'acknowledgement', 'overview'
    ]
    
    intro_keywords = ['introduction', 'intro', 'background']
    section_lower = section_title.lower().strip()
    subsection_lower = subsection_title.lower().strip() if subsection_title else ""
    
    if any(keyword in section_lower for keyword in intro_keywords) or any(keyword in subsection_lower for keyword in intro_keywords):
        return True
    
    if any(keyword in section_lower for keyword in no_citation_keywords) or any(keyword in subsection_lower for keyword in no_citation_keywords):
        return False
    
    type_lower = section_type.lower().strip()
    if any(keyword in type_lower for keyword in no_citation_keywords):
        return False
    
    return True

def extract_numeric_word_count(word_count_str):
    """Extract numeric word count from string"""
    if not word_count_str:
        return 0
    
    if isinstance(word_count_str, int):
        return word_count_str
    
    # Extract numbers from word count string
    numbers = re.findall(r'\d+', str(word_count_str))
    if numbers:
        return int(numbers[0])
    
    return 0

def create_critical_analysis_prompt_enhancement(base_prompt, section_type, word_count):
    """Enhance prompts with critical analysis requirements"""
    
    critical_analysis_addition = f"""

CRITICAL ANALYSIS REQUIREMENTS (MANDATORY):
Apply Level 2-3 Critical Analysis throughout this section:

CRITICAL ANALYSIS FRAMEWORK:
1. Use analytical verbs: evaluate, compare, synthesise, critique, argue, assess
2. Follow C-E-I structure: Claim → Evidence → Insight/Implication
3. Include methodological limitations and implications
4. Use hedged academic tone (may, appears to, suggests that)
5. Connect multiple sources/perspectives
6. Identify gaps, weaknesses, or contradictions
7. Propose future directions or improvements

CRITICAL ANALYSIS PHRASES TO INTEGRATE:
- "However, this approach suffers from certain limitations..."
- "Taken together, these findings suggest... yet remain inconclusive because..."
- "A critical evaluation reveals..."
- "Synthesising these perspectives indicates..."
- "This evidence, while compelling, must be interpreted cautiously due to..."
- "Future research should address..."

DEPTH REQUIREMENTS:
- Move beyond description to analysis and evaluation
- Weigh evidence quality and reliability
- Expose gaps in current understanding
- Judge significance and implications
- Connect to broader theoretical frameworks

Word Count: {word_count} words (maintain exact count while incorporating critical analysis)
"""
    
    return base_prompt + critical_analysis_addition


def create_subsection_generation_prompt(assignment_title, assignment_type, total_word_count, 
                                       section_title, section_type, subsection_id, subsection_title,
                                       target_word_count, numeric_word_count, content_requirements,
                                       theories, readings, learning_outcomes, rubric_criteria,
                                       specific_instructions, needs_citations, assignment_brief,
                                       module_materials):
    """Create detailed prompt for subsection generation"""
    
    citation_instruction = "Use appropriate in-text citations (e.g., Smith, 2023; Jones et al., 2022)" if needs_citations else "DO NOT include any citations in this subsection"
    
    prompt = f"""Generate high-quality academic content for this SPECIFIC SUBSECTION:

ASSIGNMENT CONTEXT:
- Assignment: {assignment_title}
- Assignment Type: {assignment_type}
- Total Word Count: {total_word_count}

SECTION CONTEXT:
- Main Section: {section_title}
- Section Type: {section_type}

SUBSECTION DETAILS:
- Subsection ID: {subsection_id}
- Subsection Title: {subsection_title}
- Target Word Count: {target_word_count}
- Content Requirements: {content_requirements}

CRITICAL WORD COUNT REQUIREMENT:
- You MUST write EXACTLY {target_word_count} (±5 words maximum)
- This is MANDATORY for THIS SUBSECTION ONLY
- Count every word carefully for this specific subsection
- Use substantive, detailed content to reach the required word count
- Never compromise on word count - expand with relevant details, examples, and analysis

CONTENT FOCUS:
- Focus ONLY on this specific subsection: "{subsection_title}"
- Address the specific requirements: {content_requirements}
- This is a SUBSECTION, not a complete section
- Do NOT write content for other subsections
- Do NOT include subsection conclusions or summaries

CRITICAL CONTENT STRUCTURE RULES:
- DO NOT include any conclusion or concluding paragraph
- DO NOT write "In conclusion", "To conclude", "To summarize", etc.
- DO NOT include any reference list or bibliography
- DO NOT write "References:" or "Bibliography:" anywhere
- {citation_instruction}
- Focus ONLY on the specific subsection content requested
- End with substantive content, not conclusions or summaries

SPECIFIC REQUIREMENTS FOR THIS SUBSECTION:"""
    # Add section interconnection
    if assignment_brief and module_materials:
        prompt += f"""

    ENHANCED INTEGRATION REQUIREMENTS:
    - Thoroughly integrate concepts from module materials: {module_materials[:500]}...
    - Ensure this subsection connects with the overall assignment flow
    - Reference relevant theories and frameworks from module materials
    - Maintain academic coherence with assignment brief requirements
    - Use transitional phrases to connect with other sections where appropriate"""


    if theories:
        prompt += f"\n- Required Theories/Models: {', '.join(theories)}"
    
    if readings:
        prompt += f"\n- Required Readings: {', '.join(readings)}"
    
    if learning_outcomes:
        prompt += f"\n- Learning Outcomes to Address: {', '.join(learning_outcomes)}"
    
    if rubric_criteria:
        prompt += f"\n- Rubric Criteria: {', '.join(rubric_criteria)}"
    
    if specific_instructions:
        prompt += f"\n- Specific Instructions: {', '.join(specific_instructions)}"

    prompt += f"""

ASSIGNMENT BRIEF:
{assignment_brief}

MODULE MATERIALS:
{module_materials}

ENHANCED GENERATION REQUIREMENTS:
1. Write complete, professional academic content for THIS SUBSECTION ONLY
2. Follow EXACT word count requirements ({target_word_count}) - this is MANDATORY
3. Include relevant theories and academic concepts from module materials
4. Use proper academic writing style with {'in-text citations' if needs_citations else 'NO citations'}
5. Include relevant examples and critical analysis to reach word count
6. NO reference list or bibliography
7. NO concluding paragraphs or summary statements
8. End with substantive content, not conclusions
9. Integrate module materials effectively
10. Meet university-level academic standards
11. Use detailed explanations, examples, and analysis to fill the word count
12. Focus specifically on: {content_requirements}

SUBSECTION-SPECIFIC STRATEGY:
- This is subsection "{subsection_title}" under section "{section_title}"
- Write ONLY for this subsection, not the entire section
- Address specific requirements: {content_requirements}
- Use {numeric_word_count} words to thoroughly cover this subsection's scope
- Provide depth and detail appropriate for this subsection
- Include specific examples and detailed analysis relevant to this subsection

IMPORTANT FINAL INSTRUCTIONS:
- Generate ONLY the content for subsection '{subsection_title}'
- Write exactly {target_word_count} (this is non-negotiable)
- {'Use in-text citations but NO reference list' if needs_citations else 'Do NOT include any citations'}
- Focus on this subsection's specific requirements
- DO NOT include any conclusion, summary, or concluding remarks
- End the subsection with substantive content, not conclusions
- Ensure professional academic writing throughout
- Make every word count towards reaching the target {target_word_count}"""
    
    # Apply critical analysis enhancement
    prompt = create_critical_analysis_prompt_enhancement(prompt, section_type, target_word_count)

    return prompt

def create_section_generation_prompt(assignment_title, assignment_type, total_word_count,
                                   section_id, section_title, section_word_count, numeric_word_count,
                                   section_type, needs_citations, assignment_brief, module_materials):
    """Create prompt for section without subsections"""
    
    citation_instruction = "Use appropriate in-text citations (e.g., Smith, 2023; Jones et al., 2022)" if needs_citations else "DO NOT include any citations in this section"
    
    prompt = f"""Generate high-quality academic content for this section:

ASSIGNMENT CONTEXT:
- Assignment: {assignment_title}
- Assignment Type: {assignment_type}
- Total Word Count: {total_word_count}

SECTION DETAILS:
- Section: {section_title}
- Section Type: {section_type}
- Target Word Count: {section_word_count}
- Section ID: {section_id}

CRITICAL WORD COUNT REQUIREMENT:
- You MUST write EXACTLY {section_word_count} (±5 words maximum)
- This is MANDATORY - count every word carefully
- Use substantive, detailed content to reach the required word count
- Never compromise on word count - expand with relevant details, examples, and analysis

CRITICAL CONTENT STRUCTURE RULES:
- DO NOT include any conclusion or concluding paragraph
- DO NOT write "In conclusion", "To conclude", "To summarize", etc.
- DO NOT include any reference list or bibliography
- DO NOT write "References:" or "Bibliography:" anywhere
- {citation_instruction}
- Focus ONLY on the specific section content requested
- End with substantive content, not conclusions or summaries

ASSIGNMENT BRIEF:
{assignment_brief}

MODULE MATERIALS:
{module_materials}

ENHANCED GENERATION REQUIREMENTS:
1. Write complete, professional academic content for this section only
2. Follow EXACT word count requirements ({section_word_count}) - this is MANDATORY
3. Include relevant theories and academic concepts from module materials
4. Use proper academic writing style with {'in-text citations' if needs_citations else 'NO citations'}
5. Include relevant examples and critical analysis to reach word count
6. NO reference list or bibliography in this section
7. NO concluding paragraphs or summary statements
8. End with substantive content, not conclusions
9. Integrate module materials effectively
10. Meet university-level academic standards
11. Use detailed explanations, examples, and analysis to fill the word count

IMPORTANT FINAL INSTRUCTIONS:
- Generate ONLY the content for the '{section_title}' section
- Write exactly {section_word_count} words (this is non-negotiable)
- {'Use in-text citations but NO reference list' if needs_citations else 'Do NOT include any citations'}
- Focus on this section's specific requirements from the brief
- DO NOT include any conclusion, summary, or concluding remarks
- End the section with substantive content, not conclusions
- Ensure professional academic writing throughout
- Make every word count towards reaching the target {section_word_count} words"""

    return prompt

def create_advanced_subsection_prompts(structure_data, assignment_brief, module_materials, assignment_id, output_dir="generated_assignments"):
    """Create advanced prompts for SUBSECTIONS with precise word count targeting"""
    
    try:
        # Extract assignment context
        assignment_title = structure_data.get("assignment_title", "Assignment")
        assignment_type = structure_data.get("assignment_type", "Academic Assignment")
        total_word_count = structure_data.get("total_word_count", "")
        sections = structure_data.get("assignment_sections", [])
        
        subsection_prompts = []
        
        logger.info(f"🔄 Creating SUBSECTION-WISE prompts for {assignment_title}...")
        
        for section in sections:
            section_id = section.get("section_id", 0)
            section_title = section.get("section_title", "")
            section_word_count = section.get("word_count", "")
            section_type = section.get("section_type", "")
            subsections = section.get("subsections", [])
            
            # Skip references section
            if (section_type.lower() in ['references', 'reference', 'bibliography'] or 
                'reference' in section_title.lower() or 
                'bibliography' in section_title.lower()):
                continue
            
            # If section has subsections, generate for each subsection separately
            if subsections and len(subsections) > 0:
                logger.info(f"📋 Section '{section_title}' has {len(subsections)} subsections - generating subsection-wise")
                
                for sub_idx, subsection in enumerate(subsections, 1):
                    subsection_id = subsection.get("subsection_id", f"{section_id}.{sub_idx}")
                    subsection_title = subsection.get("subsection_title", f"Subsection {sub_idx}")
                    subsection_word_count = subsection.get("word_count", "")
                    content_requirements = subsection.get("content_requirements", "")
                    theories = subsection.get("theories_models", [])
                    readings = subsection.get("readings", [])
                    learning_outcomes = subsection.get("learning_outcomes", [])
                    rubric_criteria = subsection.get("rubric_criteria", [])
                    specific_instructions = subsection.get("specific_instructions", [])
                    
                    # Calculate word count for subsection
                    if subsection_word_count:
                        target_word_count = subsection_word_count
                    elif section_word_count and len(subsections) > 0:
                        # Distribute section word count among subsections
                        section_numeric = extract_numeric_word_count(section_word_count)
                        if section_numeric > 0:
                            words_per_subsection = section_numeric // len(subsections)
                            target_word_count = f"{words_per_subsection} words"
                        else:
                            target_word_count = "300 words"  # Default
                    else:
                        target_word_count = "300 words"  # Default
                    
                    # Extract numeric word count for better processing
                    numeric_word_count = extract_numeric_word_count(target_word_count)
                    
                    # Determine if this subsection needs citations
                    needs_citations = determine_if_needs_citations_subsection(
                        section_title, subsection_title, section_type
                    )
                    
                    # Create focused prompt for this specific subsection
                    prompt = create_subsection_generation_prompt(
                        assignment_title=assignment_title,
                        assignment_type=assignment_type,
                        total_word_count=total_word_count,
                        section_title=section_title,
                        section_type=section_type,
                        subsection_id=subsection_id,
                        subsection_title=subsection_title,
                        target_word_count=target_word_count,
                        numeric_word_count=numeric_word_count,
                        content_requirements=content_requirements,
                        theories=theories,
                        readings=readings,
                        learning_outcomes=learning_outcomes,
                        rubric_criteria=rubric_criteria,
                        specific_instructions=specific_instructions,
                        needs_citations=needs_citations,
                        assignment_brief=assignment_brief,
                        module_materials=module_materials
                    )
                    
                    subsection_prompts.append({
                        "section_id": section_id,
                        "section_title": section_title,
                        "subsection_id": subsection_id,
                        "subsection_title": subsection_title,
                        "word_count": target_word_count,
                        "numeric_word_count": numeric_word_count,
                        "section_type": section_type,
                        "assignment_type": assignment_type,
                        "needs_citations": needs_citations,
                        "content_requirements": content_requirements,
                        "is_subsection": True,
                        "prompt": prompt
                    })
                    
            else:
                # No subsections, treat as regular section
                logger.info(f"📋 Section '{section_title}' has no subsections - generating as single section")
                
                # Extract numeric word count for better processing
                numeric_word_count = extract_numeric_word_count(section_word_count)
                
                # Determine if this section needs citations
                needs_citations = determine_if_needs_citations(section_title, section_type)
                
                # Create prompt for the entire section
                prompt = create_section_generation_prompt(
                    assignment_title=assignment_title,
                    assignment_type=assignment_type,
                    total_word_count=total_word_count,
                    section_id=section_id,
                    section_title=section_title,
                    section_word_count=section_word_count,
                    numeric_word_count=numeric_word_count,
                    section_type=section_type,
                    needs_citations=needs_citations,
                    assignment_brief=assignment_brief,
                    module_materials=module_materials
                )
                
                subsection_prompts.append({
                    "section_id": section_id,
                    "section_title": section_title,
                    "subsection_id": None,
                    "subsection_title": None,
                    "word_count": section_word_count,
                    "numeric_word_count": numeric_word_count,
                    "section_type": section_type,
                    "assignment_type": assignment_type,
                    "needs_citations": needs_citations,
                    "content_requirements": "",
                    "is_subsection": False,
                    "prompt": prompt
                })
        
        # Save prompts
        assignment_dir = os.path.join(output_dir, f"assignment_{assignment_id}")
        prompts_file = os.path.join(assignment_dir, "subsection_prompts.json")
        with open(prompts_file, "w", encoding="utf-8") as f:
            json.dump(subsection_prompts, f, indent=2, ensure_ascii=False)
        
        logger.info(f"✅ Created {len(subsection_prompts)} SUBSECTION-WISE prompts")
        logger.info(f"📁 Saved to '{prompts_file}'")
        
        return subsection_prompts
        
    except Exception as e:
        logger.error(f"Error creating subsection prompts: {e}")
        return None

def clean_generated_content(content):
    """Clean generated content removing unwanted elements"""
    
    # Enhanced content cleaning
    content_lines = content.split('\n')
    cleaned_lines = []
    
    for line in content_lines:
        line_lower = line.lower().strip()
        # Skip unwanted lines
        if (line_lower.startswith('references') or 
            line_lower.startswith('bibliography') or 
            line_lower.startswith('in conclusion') or
            line_lower.startswith('to conclude') or
            line_lower.startswith('conclusion') or
            line_lower.startswith('word count:')):
            break
        cleaned_lines.append(line)
    
    content = '\n'.join(cleaned_lines).strip()
    
    # Remove concluding sentences more aggressively
    sentences = content.split('.')
    filtered_sentences = []
    
    for sentence in sentences:
        sentence_lower = sentence.lower().strip()
        if (sentence_lower.startswith('in conclusion') or
            sentence_lower.startswith('to conclude') or
            sentence_lower.startswith('overall') or
            sentence_lower.startswith('in summary') or
            'in conclusion' in sentence_lower or
            'to conclude' in sentence_lower):
            continue
        filtered_sentences.append(sentence)
    
    content = '.'.join(filtered_sentences).strip()
    if content and not content.endswith('.'):
        content += '.'
    
    return content

def generate_single_subsection_advanced(subsection_prompt_data, assignment_id, output_dir="generated_assignments"):
    """Enhanced subsection generation with precise word count compliance"""
    
    section_id = subsection_prompt_data.get("section_id")
    section_title = subsection_prompt_data.get("section_title")
    subsection_id = subsection_prompt_data.get("subsection_id")
    subsection_title = subsection_prompt_data.get("subsection_title")
    is_subsection = subsection_prompt_data.get("is_subsection", False)
    section_type = subsection_prompt_data.get("section_type", "")
    assignment_type = subsection_prompt_data.get("assignment_type", "")
    word_count = subsection_prompt_data.get("word_count", "")
    numeric_word_count = subsection_prompt_data.get("numeric_word_count", 0)
    needs_citations = subsection_prompt_data.get("needs_citations", True)
    prompt = subsection_prompt_data.get("prompt")
    
    assignment_dir = os.path.join(output_dir, f"assignment_{assignment_id}", "subsections")
    os.makedirs(assignment_dir, exist_ok=True)
    
    if is_subsection:
        generation_target = f"{section_title} > {subsection_title}"
        logger.info(f"🔄 Generating SUBSECTION: {generation_target} - Target: {word_count} words - Citations: {'Yes' if needs_citations else 'No'}")
    else:
        generation_target = section_title
        logger.info(f"🔄 Generating SECTION: {generation_target} - Target: {word_count} words - Citations: {'Yes' if needs_citations else 'No'}")
    
    try:
        # Enhanced system message with better token allocation
        citation_instruction = "Use ONLY in-text citations (e.g., Smith, 2023; Jones et al., 2022)" if needs_citations else "DO NOT include any citations"
        content_type = "subsection" if is_subsection else "section"
        
        system_msg = f"""You are an expert academic writer specializing in precise word count compliance for {content_type} generation.

        CRITICAL REQUIREMENTS:
        1. Write EXACTLY {word_count} words (±3 words maximum)
        2. {citation_instruction}
        3. DO NOT include any reference list or bibliography
        4. DO NOT write "References:" anywhere in the content
        5. DO NOT include any conclusion or concluding paragraph
        6. DO NOT write "In conclusion", "To conclude", "To summarize", etc.
        7. Focus ONLY on the specific {content_type} content requested
        8. End with substantive content, not conclusions or summaries
        9. Integrate module materials effectively
        10. Use professional academic writing style throughout
        11. WORD COUNT IS MANDATORY - expand content to meet exact requirement

        {content_type.upper()} GENERATION STRATEGY:
        - Use detailed explanations and comprehensive analysis
        - Include specific examples and case studies
        - Provide thorough background and context
        - Discuss implications and applications in detail
        - Use multiple well-developed paragraphs for longer {content_type}s
        """
        
        # Determine appropriate max_tokens based on word count
        if numeric_word_count > 800:
            max_tokens = 6000
        elif numeric_word_count > 500:
            max_tokens = 4000
        elif numeric_word_count > 300:
            max_tokens = 3000
        else:
            max_tokens = 2000
        
        # Multiple attempts for better word count compliance
        best_content = None
        best_word_difference = float('inf')
        
        for attempt in range(3):
            try:
                response = openai_client.chat.completions.create(
                    model="gpt-4o",
                    messages=[
                        {"role": "system", "content": system_msg},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.2 + (attempt * 0.1),  # Slightly increase temperature for variety
                    max_tokens=max_tokens
                )
                
                content = response.choices[0].message.content
                
                # Enhanced content cleaning
                content = clean_generated_content(content)
                
                # Verify word count
                actual_word_count = len(content.split())
                target_words = numeric_word_count if numeric_word_count > 0 else 0
                
                if target_words > 0:
                    word_difference = abs(actual_word_count - target_words)
                    logger.info(f"   Attempt {attempt + 1}: Target={target_words}, Actual={actual_word_count}, Difference={word_difference}")
                    
                    # Track best attempt
                    if word_difference < best_word_difference:
                        best_word_difference = word_difference
                        best_content = content
                    
                    # If word count is acceptable, break
                    if word_difference <= 10:
                        break
                    
                    # If this is the last attempt, use best content
                    if attempt == 2:
                        content = best_content
                        logger.warning(f"   Using best content with {best_word_difference} word difference after 3 attempts")
                        break
                        
                else:
                    # No specific word count target, use first result
                    break
                    
            except Exception as api_error:
                logger.error(f"   API error on attempt {attempt + 1}: {api_error}")
                if attempt == 2:  # Last attempt
                    raise api_error
                time.sleep(1)  # Wait before retry
        
        # Create filename
        if is_subsection:
            clean_section_title = re.sub(r'[^\w\s-]', '', section_title).strip()
            clean_section_title = re.sub(r'[-\s]+', '_', clean_section_title)
            clean_subsection_title = re.sub(r'[^\w\s-]', '', subsection_title).strip()
            clean_subsection_title = re.sub(r'[-\s]+', '_', clean_subsection_title)
            filename = os.path.join(assignment_dir, f"section_{section_id}_{clean_section_title}_sub_{subsection_id}_{clean_subsection_title}.md")
        else:
            clean_title = re.sub(r'[^\w\s-]', '', section_title).strip()
            clean_title = re.sub(r'[-\s]+', '_', clean_title)
            filename = os.path.join(assignment_dir, f"section_{section_id}_{clean_title}.md")
        
        with open(filename, "w", encoding="utf-8") as f:
            if is_subsection:
                f.write(f"# {section_title}\n\n## {subsection_title}\n\n")
            else:
                f.write(f"# {section_title}\n\n")
            f.write(content)
        
        final_word_count = len(content.split())
        logger.info(f"✅ Generated: {filename} ({final_word_count} words)")
        return filename
        
    except Exception as e:
        logger.error(f"Error generating {generation_target}: {e}")
        return None

def generate_all_assignment_subsections_advanced(subsection_prompts, assignment_id):
    """Generate all subsections with enhanced progress tracking"""
    
    try:
        logger.info(f"🚀 Starting SUBSECTION-WISE generation for {len(subsection_prompts)} parts")
        logger.info("=" * 80)
        
        generated_files = []
        
        for i, subsection_data in enumerate(subsection_prompts, 1):
            logger.info(f"\n📝 Processing part {i}/{len(subsection_prompts)}")
            
            filename = generate_single_subsection_advanced(subsection_data, assignment_id)
            if filename:
                generated_files.append(filename)
            
            # Rate limiting
            time.sleep(2)
        
        logger.info(f"\n🎯 SUBSECTION-WISE GENERATION COMPLETE!")
        logger.info(f"✅ Successfully generated {len(generated_files)} parts")
        
        return generated_files
        
    except Exception as e:
        logger.error(f"Error in subsection-wise generation: {e}")
        return None

def add_professional_title_page_advanced(doc, assignment_title, assignment_type, total_word_count):
    """Add professional title page with university placeholders"""
    
    # University placeholder
    uni_para = doc.add_paragraph()
    uni_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    uni_run = uni_para.add_run("[UNIVERSITY NAME]")
    uni_run.font.size = Pt(14)
    uni_run.font.name = 'Times New Roman'
    uni_run.bold = True
    
    # Department
    dept_para = doc.add_paragraph()
    dept_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    dept_run = dept_para.add_run("[Department/Faculty Name]")
    dept_run.font.size = Pt(12)
    dept_run.font.name = 'Times New Roman'
    
    # Spacing
    for _ in range(4):
        doc.add_paragraph()
    
    # Assignment title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run(assignment_title.upper())
    title_run.font.size = Pt(18)
    title_run.bold = True
    title_run.font.name = 'Times New Roman'
    
    doc.add_paragraph()
    
    # Assignment type
    type_para = doc.add_paragraph()
    type_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    type_run = type_para.add_run(assignment_type)
    type_run.font.size = Pt(14)
    type_run.font.name = 'Times New Roman'
    
    # Word count
    if total_word_count and total_word_count != "Not specified":
        doc.add_paragraph()
        wordcount_para = doc.add_paragraph()
        wordcount_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        wordcount_run = wordcount_para.add_run(f"Word Count: {total_word_count}")
        wordcount_run.font.size = Pt(12)
        wordcount_run.font.name = 'Times New Roman'
    
    # Spacing
    for _ in range(6):
        doc.add_paragraph()
    
    # Student details
    details_para = doc.add_paragraph()
    details_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    details_text = """Submitted by:
    [Student Name]
    Student ID: [Student ID]
    Course: [Course Name]
    Module: [Module Code]

    Submitted to:
    [Instructor Name]
    [Title]

    Submission Date: [Date]"""
    
    details_run = details_para.add_run(details_text)
    details_run.font.size = Pt(12)
    details_run.font.name = 'Times New Roman'
    details_para.paragraph_format.line_spacing = 1.5

def add_professional_toc_advanced(doc, structure_data):
    """Add professional table of contents with proper formatting"""
    
    # TOC Title
    toc_title_para = doc.add_paragraph()
    toc_title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    toc_title_run = toc_title_para.add_run("TABLE OF CONTENTS")
    toc_title_run.font.size = Pt(16)
    toc_title_run.bold = True
    toc_title_run.font.name = 'Times New Roman'
    
    doc.add_paragraph()
    
    # Generate TOC entries
    sections = structure_data.get("assignment_sections", [])
    page_num = 3
    section_counter = 1

    for section in sections:
        section_title = section.get("section_title", f"Section {section_counter}")
        section_type = section.get("section_type", "")
        subsections = section.get("subsections", [])
        
        # Skip references section in main TOC
        if (section_type.lower() in ['references', 'reference', 'bibliography'] or 
            'reference' in section_title.lower()):
            continue
        
        toc_para = doc.add_paragraph()
        toc_para.paragraph_format.left_indent = Inches(0)
        
        section_text = f"{section_counter}. {section_title}"
        toc_text_run = toc_para.add_run(section_text)
        toc_text_run.font.name = 'Times New Roman'
        toc_text_run.font.size = Pt(12)
        toc_text_run.bold = True
        
        dots_needed = 60 - len(section_text)
        if dots_needed > 0:
            dots = "." * dots_needed
            toc_para.add_run(f" {dots} ")
        
        page_run = toc_para.add_run(str(page_num))
        page_run.font.name = 'Times New Roman'
        page_run.font.size = Pt(12)
        page_run.bold = True
        
        if subsections and len(subsections) > 0:
            for sub_idx, subsection in enumerate(subsections, 1):
                subsection_title = subsection.get("subsection_title", f"Subsection {sub_idx}")
                
                sub_toc_para = doc.add_paragraph()
                sub_toc_para.paragraph_format.left_indent = Inches(0.3)
                
                sub_text = f"{section_counter}.{sub_idx} {subsection_title}"
                sub_toc_run = sub_toc_para.add_run(sub_text)
                sub_toc_run.font.name = 'Times New Roman'
                sub_toc_run.font.size = Pt(11)
                
                sub_dots_needed = 55 - len(sub_text)
                if sub_dots_needed > 0:
                    sub_dots = "." * sub_dots_needed
                    sub_toc_para.add_run(f" {sub_dots} ")
                
                sub_page_run = sub_toc_para.add_run(str(page_num))
                sub_page_run.font.name = 'Times New Roman'
                sub_page_run.font.size = Pt(11)
        
        section_counter += 1
        page_num += 1
    
    # Add References to TOC
    ref_para = doc.add_paragraph()
    ref_para.paragraph_format.left_indent = Inches(0)
    
    ref_text = "References"
    ref_text_run = ref_para.add_run(ref_text)
    ref_text_run.font.name = 'Times New Roman'
    ref_text_run.font.size = Pt(12)
    ref_text_run.bold = True
    
    ref_dots_needed = 70 - len(ref_text)
    if ref_dots_needed > 0:
        ref_dots = "." * ref_dots_needed
        ref_para.add_run(f" {ref_dots} ")
    
    ref_page_run = ref_para.add_run(str(page_num))
    ref_page_run.font.name = 'Times New Roman'
    ref_page_run.font.size = Pt(12)
    ref_page_run.bold = True

def add_section_heading_to_doc(doc, section_title, section_num):
    """Add main section heading to document"""
    
    heading_para = doc.add_paragraph()
    heading_para.paragraph_format.space_before = Pt(12)
    heading_para.paragraph_format.space_after = Pt(6)
    
    heading_run = heading_para.add_run(f"{section_num}. {section_title.upper()}")
    heading_run.font.size = Pt(14)
    heading_run.bold = True
    heading_run.font.name = 'Times New Roman'
    
    # Add underline
    underline_para = doc.add_paragraph()
    underline_run = underline_para.add_run("_" * 60)
    underline_run.font.size = Pt(10)
    underline_para.paragraph_format.space_after = Pt(12)

def add_subsection_to_doc(doc, subsection_title, content, sub_idx):
    """Add subsection content to document"""
    
    # Remove markdown headers from content
    content_lines = content.split('\n')
    cleaned_content_lines = []
    
    for line in content_lines:
        if line.startswith('#'):
            continue  # Skip markdown headers
        cleaned_content_lines.append(line)
    
    content = '\n'.join(cleaned_content_lines).strip()
    
    # Add subsection heading
    sub_heading_para = doc.add_paragraph()
    sub_heading_para.paragraph_format.space_before = Pt(12)
    sub_heading_para.paragraph_format.space_after = Pt(6)
    sub_heading_para.paragraph_format.left_indent = Inches(0.25)
    
    sub_heading_run = sub_heading_para.add_run(f"{sub_idx}. {subsection_title}")
    sub_heading_run.font.size = Pt(12)
    sub_heading_run.bold = True
    sub_heading_run.font.name = 'Times New Roman'
    
    # Add subsection content
    add_content_paragraphs_to_doc(doc, content)

def add_section_content_to_doc(doc, content):
    """Add section content to document (for sections without subsections)"""
    
    # Remove markdown headers from content
    content_lines = content.split('\n')
    cleaned_content_lines = []
    
    for line in content_lines:
        if line.startswith('#'):
            continue  # Skip markdown headers
        cleaned_content_lines.append(line)
    
    content = '\n'.join(cleaned_content_lines).strip()
    
    # Add content
    add_content_paragraphs_to_doc(doc, content)

def add_content_paragraphs_to_doc(doc, content):
    """Add content paragraphs to document with proper formatting"""
    
    paragraphs = content.split('\n\n')
    
    for para_text in paragraphs:
        para_text = para_text.strip()
        if not para_text:
            continue
        
        # Regular content paragraph
        content_para = doc.add_paragraph()
        content_para.paragraph_format.space_after = Pt(6)
        content_para.paragraph_format.line_spacing = 1.5
        content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        content_para.paragraph_format.first_line_indent = Inches(0.5)
        
        content_run = content_para.add_run(para_text)
        content_run.font.name = 'Times New Roman'
        content_run.font.size = Pt(11)

def find_subsection_file(subsections_dir, section_id, section_title, subsection_id, subsection_title):
    """Find subsection file dynamically"""
    
    if not os.path.exists(subsections_dir):
        return None
    
    clean_section_title = re.sub(r'[^\w\s-]', '', section_title).strip()
    clean_section_title = re.sub(r'[-\s]+', '_', clean_section_title)
    clean_subsection_title = re.sub(r'[^\w\s-]', '', subsection_title).strip()
    clean_subsection_title = re.sub(r'[-\s]+', '_', clean_subsection_title)
    
    possible_names = [
        f"section_{section_id}_{clean_section_title}_sub_{subsection_id}_{clean_subsection_title}.md",
        f"section_{section_id}_{clean_section_title}_sub_{clean_subsection_title}.md",
        f"{clean_section_title}_sub_{clean_subsection_title}.md"
    ]
    
    for filename in possible_names:
        filepath = os.path.join(subsections_dir, filename)
        if os.path.exists(filepath):
            return filepath
    
    return None

def find_section_file_subsections(subsections_dir, section_id, section_title):
    """Find section file in subsections directory"""
    
    if not os.path.exists(subsections_dir):
        return None
    
    clean_title = re.sub(r'[^\w\s-]', '', section_title).strip()
    clean_title = re.sub(r'[-\s]+', '_', clean_title)
    
    possible_names = [
        f"section_{section_id}_{clean_title}.md",
        f"{clean_title}.md"
    ]
    
    for filename in possible_names:
        filepath = os.path.join(subsections_dir, filename)
        if os.path.exists(filepath):
            return filepath
    
    return None

def extract_citations_from_subsections(subsections_dir):
    """Extract all citations from generated subsection files"""
    
    citations = set()  # Use set to avoid duplicates
    
    if not os.path.exists(subsections_dir):
        return list(citations)
    
    # Advanced citation patterns
    citation_patterns = [
        r'\(([A-Za-z]+(?:\s+et\s+al\.?)?,?\s+\d{4}[a-z]?)\)',  # (Smith, 2023) or (Smith et al., 2023)
        r'\(([A-Za-z]+\s+&\s+[A-Za-z]+,?\s+\d{4}[a-z]?)\)',   # (Smith & Jones, 2023)
        r'\(([A-Z]+,?\s+\d{4}[a-z]?)\)',                       # (CIPD, 2023)
        r'\(([A-Za-z]+(?:\s+et\s+al\.?)?,?\s+\d{4}[a-z]?;\s*[A-Za-z]+(?:\s+et\s+al\.?)?,?\s+\d{4}[a-z]?)\)',  # Multiple citations
    ]
    
    try:
        # Read all markdown files in subsections directory
        for filename in os.listdir(subsections_dir):
            if filename.endswith('.md'):
                filepath = os.path.join(subsections_dir, filename)
                with open(filepath, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                # Extract citations using regex patterns
                for pattern in citation_patterns:
                    matches = re.findall(pattern, content)
                    for match in matches:
                        # Handle multiple citations in one parenthesis
                        if ';' in match:
                            individual_citations = match.split(';')
                            for citation in individual_citations:
                                clean_citation = citation.strip().replace(',', '').strip()
                                if clean_citation and len(clean_citation) > 3:
                                    citations.add(clean_citation)
                        else:
                            # Clean up the citation
                            clean_citation = match.strip().replace(',', '').strip()
                            if clean_citation and len(clean_citation) > 3:  # Basic validation
                                citations.add(clean_citation)
    
    except Exception as e:
        logger.error(f"Error reading subsection files: {e}")
    
    return sorted(list(citations))

def generate_full_references_advanced(citations):
    """Generate full references from extracted citations using advanced AI"""
    if not citations:
        return []
    # Create prompt for generating full references
    citations_text = "\n".join([f"- {citation}" for citation in citations])
    prompt = f"""Convert these in-text citations into full Harvard-style references. 

IMPORTANT RULES:
1. Generate ONLY references for the citations provided below
2. Use proper Harvard referencing format
3. ALL REFERENCES MUST BE FROM 2020-2025 ONLY (last 5 years, inclusive)
4. Make realistic but plausible academic references from the specified timeframe
5. Include proper publication details (publisher, edition, etc.)
6. Sort alphabetically by author surname
7. Do NOT add any extra references not mentioned in the citations
8. Ensure all publication years are between 2020-2025 (inclusive)

CRITICAL: Every reference must have a publication year between 2020-2025 (inclusive). If a citation cannot be matched to a reference in this range, SKIP IT.

In-text citations found in the assignment:
{citations_text}

Generate the complete reference list in Harvard format with years 2020-2025 only. Do NOT include any reference outside this range. If a reference cannot be made with a year in this range, SKIP IT."""
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an expert in academic referencing. Generate accurate Harvard-style references only for the citations provided."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=2000
        )
        references_text = response.choices[0].message.content.strip()
        # Split into individual references and clean them
        references = []
        for line in references_text.split('\n'):
            line = line.strip()
            if line and not line.startswith('#') and len(line) > 10:
                # Remove any bullet points or numbering
                clean_ref = re.sub(r'^\d+\.\s*', '', line)
                clean_ref = re.sub(r'^[-•]\s*', '', clean_ref)
                clean_ref = clean_ref.strip()
                # Post-process: Only keep references with year 2020-2025
                year_match = re.search(r'(20[2-2][0-9]|2025)', clean_ref)
                if year_match:
                    year = int(year_match.group(0))
                    if 2020 <= year <= 2025:
                        references.append(clean_ref)
        return references
    except Exception as e:
        logger.error(f"Error generating references: {e}")
        return []

def add_professional_references_section_from_subsections(doc, subsections_dir):
    """Add professional references section with citations from subsection files"""
    
    # References heading
    ref_heading_para = doc.add_paragraph()
    ref_heading_para.paragraph_format.space_before = Pt(12)
    ref_heading_para.paragraph_format.space_after = Pt(6)
    
    ref_heading_run = ref_heading_para.add_run("REFERENCES")
    ref_heading_run.font.size = Pt(14)
    ref_heading_run.bold = True
    ref_heading_run.font.name = 'Times New Roman'
    
    # Add underline
    underline_para = doc.add_paragraph()
    underline_run = underline_para.add_run("_" * 60)
    underline_run.font.size = Pt(10)
    underline_para.paragraph_format.space_after = Pt(12)
    
    # Extract citations from all generated subsection files
    try:
        citations = extract_citations_from_subsections(subsections_dir)
        
        if citations:
            # Generate full references for extracted citations
            references = generate_full_references_advanced(citations)
            
            logger.info(f"✅ Generated {len(references)} references from {len(citations)} citations")
            
            for ref in references:
                ref_para = doc.add_paragraph()
                ref_para.paragraph_format.left_indent = Inches(0.5)
                ref_para.paragraph_format.first_line_indent = Inches(-0.5)  # Hanging indent
                ref_para.paragraph_format.space_after = Pt(6)
                ref_para.paragraph_format.line_spacing = 1.15
                
                ref_run = ref_para.add_run(ref)
                ref_run.font.name = 'Times New Roman'
                ref_run.font.size = Pt(11)
        else:
            # If no citations found, add a note
            note_para = doc.add_paragraph()
            note_para.paragraph_format.first_line_indent = Inches(0.5)
            note_run = note_para.add_run("[References will be added based on citations used in the assignment content]")
            note_run.italic = True
            note_run.font.name = 'Times New Roman'
            note_run.font.size = Pt(11)
            
    except Exception as e:
        logger.warning(f"Could not extract citations automatically: {e}")
        # Fallback to placeholder
        note_para = doc.add_paragraph()
        note_para.paragraph_format.first_line_indent = Inches(0.5)
        note_run = note_para.add_run("[References will be added based on citations used in the assignment content]")
        note_run.italic = True
        note_run.font.name = 'Times New Roman'
        note_run.font.size = Pt(11)

def apply_professional_formatting_advanced(doc):
    """Apply consistent professional formatting with advanced rules"""
    
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            # Ensure consistent font
            for run in paragraph.runs:
                if not run.font.name:
                    run.font.name = 'Times New Roman'
                if not run.font.size:
                    run.font.size = Pt(11)
            
            # Set paragraph spacing
            if paragraph.paragraph_format.space_after is None:
                paragraph.paragraph_format.space_after = Pt(6)
            
            # Line spacing for body text
            if (paragraph.paragraph_format.line_spacing is None and 
                not any(run.bold for run in paragraph.runs)):
                paragraph.paragraph_format.line_spacing = 1.5

def create_professional_docx_from_subsections(assignment_id, structure_data, output_dir="generated_assignments"):
    """Create professional DOCX document from subsection files"""
    
    try:
        assignment_dir = os.path.join(output_dir, f"assignment_{assignment_id}")
        subsections_dir = os.path.join(assignment_dir, "subsections")
        
        # Extract assignment details
        assignment_title = structure_data.get("assignment_title", "Assignment")
        assignment_type = structure_data.get("assignment_type", "Academic Assignment")
        total_word_count = structure_data.get("total_word_count", "")
        
        # Dynamic output filename
        clean_title = re.sub(r'[^\w\s-]', '', assignment_title).strip()
        clean_title = re.sub(r'[-\s]+', '_', clean_title)
        output_file = os.path.join(assignment_dir, f"{clean_title}_Complete_Subsections.docx")
        
        logger.info(f"📄 Creating professional DOCX from subsections for: {assignment_title}")
        
        # Create document
        doc = Document()
        
        # Set document margins and page setup
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1.25)
            section.right_margin = Inches(1.25)
            section.page_height = Inches(11.69)  # A4 height
            section.page_width = Inches(8.27)   # A4 width
        
        # Add professional title page
        add_professional_title_page_advanced(doc, assignment_title, assignment_type, total_word_count)
        doc.add_page_break()
        
        # Add professional table of contents
        add_professional_toc_advanced(doc, structure_data)
        doc.add_page_break()
        
        # Add sections and subsections
        sections_data = structure_data.get("assignment_sections", [])
        section_counter = 1
        
        for section_info in sections_data:
            section_title = section_info.get("section_title", f"Section {section_counter}")
            section_type = section_info.get("section_type", "")
            subsections = section_info.get("subsections", [])
            
            # Skip references section
            if (section_type.lower() in ['references', 'reference', 'bibliography'] or 
                'reference' in section_title.lower()):
                continue
            
            logger.info(f"   Adding: {section_title}")
            
            # Add main section heading
            add_section_heading_to_doc(doc, section_title, section_counter)
            
            if subsections and len(subsections) > 0:
                # Process subsections
                for sub_idx, subsection in enumerate(subsections, 1):
                    subsection_title = subsection.get("subsection_title", f"Subsection {sub_idx}")
                    subsection_id = subsection.get("subsection_id", f"{section_counter}.{sub_idx}")
                    
                    # Find subsection file
                    subsection_file = find_subsection_file(subsections_dir, section_counter, section_title, subsection_id, subsection_title)
                    
                    if subsection_file and os.path.exists(subsection_file):
                        with open(subsection_file, 'r', encoding='utf-8') as f:
                            content = f.read()
                        
                        add_subsection_to_doc(doc, subsection_title, content, sub_idx)
                    else:
                        logger.warning(f"   ⚠️ Warning: Subsection file not found for {subsection_title}")
            else:
                # No subsections, find section file
                section_file = find_section_file_subsections(subsections_dir, section_counter, section_title)
                
                if section_file and os.path.exists(section_file):
                    with open(section_file, 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    add_section_content_to_doc(doc, content)
                else:
                    logger.warning(f"   ⚠️ Warning: Section file not found for {section_title}")
            
            section_counter += 1
            doc.add_page_break()
        
        # Add professional references section
        logger.info("   Extracting citations and generating references...")
        add_professional_references_section_from_subsections(doc, subsections_dir)
        
        # Apply consistent formatting
        apply_professional_formatting_advanced(doc)
        
        # Save
        doc.save(output_file)
        
        # Statistics
        total_paragraphs = len(doc.paragraphs)
        estimated_words = sum(len(p.text.split()) for p in doc.paragraphs if p.text.strip())
        
        logger.info(f"\n✅ Complete assignment saved as '{output_file}'")
        logger.info(f"📊 Statistics:")
        logger.info(f"   - Paragraphs: {total_paragraphs}")
        logger.info(f"   - Estimated words: {estimated_words}")
        logger.info(f"   - Professional formatting applied")
        
        return output_file
        
    except Exception as e:
        logger.error(f"Error creating DOCX from subsections: {e}")
        return None

# ORIGINAL EXTRACTION FUNCTIONS (keeping all existing functions)

# function to store details of failed files in a csv file
def log_failed_processing(file_path, file_extension, error_message, original_filename=None):
    try:
        # Get the current working directory
        current_dir = os.getcwd()
        csv_file = os.path.join(current_dir, "failed_processing_log.csv")
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        # Use original filename if provided, otherwise extract from path
        file_name = original_filename if original_filename else (os.path.basename(file_path) if file_path else "No file")
        
        # Prepare the new record
        new_record = [timestamp, file_name, file_extension, error_message]
        
        # Read existing content if file exists
        existing_records = []
        if os.path.exists(csv_file):
            with open(csv_file, 'r', newline='', encoding='utf-8') as f:
                reader = csv.reader(f)
                existing_records = list(reader)
        
        # Write all content back to file with new record at the beginning
        with open(csv_file, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            # Write headers if file is new
            if not existing_records:
                writer.writerow(['Timestamp', 'Customer ID', 'File Name', 'File Extension', 'Error Message'])
            else:
                # Write headers from existing file
                writer.writerow(existing_records[0])
            
            # Write new record first
            writer.writerow(new_record)
            
            # Write existing records (skip header if exists)
            if existing_records:
                writer.writerows(existing_records[1:])
            
            logger.info(f"Logged failed processing to CSV: {csv_file}")
    except Exception as e:
        logger.error(f"Error writing to CSV file: {str(e)}", exc_info=True)

def extract_text_from_ppt_with_images(file_path: str) -> str:
    """Extract text from a PPT, including text in images using OCR."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            text += extract_shape_text_with_images(shape)
    return text

# function to extract text from shapes 
def extract_shape_text_with_images(shape) -> str:
    """Extract text from shapes and perform OCR on images."""
    text = ""
    if shape.has_text_frame:
        text += shape.text_frame.text + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            text += extract_shape_text_with_images(sub_shape)
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            image_bytes = image.blob
            tmp.write(image_bytes)
            tmp_path = tmp.name
        ocr_text = extract_text_from_image(tmp_path)
        text += ocr_text + "\n"
        os.remove(tmp_path)
    return text

# fucntion to extract text from images
def extract_text_from_image(file_path: str) -> str:
    """Extract text from an image file using Google Vision API."""
    try:
        # Load image into memory
        with open(file_path, 'rb') as image_file:
            content = image_file.read()

        # Use Vision API to detect text
        image = vision.Image(content=content)
        response = client.text_detection(image=image)
        texts = response.text_annotations

        if response.error.message:
            logger.error(f"Google Vision API error: {response.error.message}")
            return ""

        # Return the full text from the first text annotation (the most complete one)
        if texts:
            return texts[0].description
        else:
            return ""
    except Exception as e:
        logger.error(f"Error during OCR processing: {str(e)}", exc_info=True)
        return ""

# word file extraction 
def extract_text_from_word(word_file):
    try:
        # If the word file is a path (string), open it as binary, otherwise, handle the file object
        if isinstance(word_file, str):
            with open(word_file, 'rb') as f:
                word_file_content = BytesIO(f.read())
        else:
            word_file_content = BytesIO(word_file.read())
        
        # Load the document
        doc = Document(word_file_content)
        
        # Extract paragraphs
        text = "\n".join([para.text for para in doc.paragraphs])
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text  # Add text from each cell in the table
      
        # Optionally print to debug or analyze
        # print(text)  # For debugging purposes
        
        return text
    except Exception as e:
        logger.error(f"Error extracting text from Word document: {str(e)}", exc_info=True)
        return ""
    
# pdf to image conversion
def pdf_to_images(file_path: str, request_temp_dir: str) -> list:
    """
    Converts each page of a PDF to an image (PNG format) and saves them to disk.
    
    Args:
        pdf_path: Path to the PDF file.
        output_folder: Folder where the images will be saved.
        
    Returns:
        A list of paths to the generated image files.
    """
    # Ensure the output directory exists
    output_folder = os.path.join(request_temp_dir, f"output_images_{uuid.uuid4()}")
    os.makedirs(output_folder, exist_ok=True)
    
    # Open the PDF file with PyMuPDF
    pdf_document = fitz.open(file_path)
    image_paths = []

    # Loop through each page in the PDF
    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]
        # Render the page to an image (pixmap)
        pix = page.get_pixmap()
        # Define the path for this image. Using PNG preserves quality.
        img_path = os.path.join(output_folder, f"page_{page_number+1}.png")
        pix.save(img_path)
        image_paths.append(img_path)
        print(f"Saved page {page_number+1} as image at {img_path}")

    print(f"Converted {len(image_paths)} pages to images.")
    return image_paths

# word file to images conversion
def word_to_images(file_path: str, request_temp_dir: str) -> list:
    # Create necessary folders
    output_folder = os.path.join(request_temp_dir, f"output_images_{uuid.uuid4()}")
    os.makedirs(output_folder, exist_ok=True)
    pdf_output = os.path.join(output_folder, "converted.pdf")

    # Convert DOCX to PDF using LibreOffice
    print(f"Converting DOCX to PDF: {file_path}")
    subprocess.run([
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_folder,
        file_path
    ], check=True)

    # Find the converted PDF (same name, different extension)
    base_filename = os.path.splitext(os.path.basename(file_path))[0]
    pdf_path = os.path.join(output_folder, base_filename + ".pdf")

    # Convert PDF to images
    print(f"Converting PDF to images: {pdf_path}")
    images = convert_from_path(pdf_path)
    image_paths = []

    for i, img in enumerate(images):
        img_path = os.path.join(output_folder, f"page_{i+1}.png")
        img.save(img_path, "PNG")
        image_paths.append(img_path)
        print(f"Saved page {i+1} as image at {img_path}")

    # Clean up PDF
    if os.path.exists(pdf_path):
        os.remove(pdf_path)

    return image_paths

# convert word to images using aspose word library
def WORD_TO_IMAGES(file_path: str, request_temp_dir: str) -> list:
   
    # Ensure the output directory exists
    output_folder = os.path.join(request_temp_dir, f"output_images_{uuid.uuid4()}")
    os.makedirs(output_folder, exist_ok=True)

    # Load the Word document
    doc = aw.Document(file_path)
    image_paths = []

    # Set image save options
    save_options = aw.saving.ImageSaveOptions(aw.SaveFormat.PNG)

    # Loop through each page in the document
    for page_number in range(doc.page_count):
        # Specify the page to save
        save_options.page_set = aw.saving.PageSet(page_number)
        # Define the path for this image
        img_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
        # Save the page as an image
        doc.save(img_path, save_options)
        image_paths.append(img_path)
        print(f"Saved page {page_number + 1} as image at {img_path}")

    print(f"Converted {len(image_paths)} pages to images.")
    return image_paths

# Function for text extraction from pdf files
def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file, including OCR on images."""
    text = ""
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text += page.get_text() + "\n"
            # Extract images and perform OCR
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                    tmp.write(image_bytes)
                    tmp_path = tmp.name
                ocr_text = extract_text_from_image(tmp_path)
                text += ocr_text + "\n"
                os.remove(tmp_path)
        return text
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {str(e)}", exc_info=True)
        return ""

def extract_from_archive(file_path: str, file_type: str, temp_dir: str) -> tuple[str, str]:
    """
    Extract and process files from ZIP, RAR, or TAR archives.
    Returns a tuple of (main_text, rubric_text)
    """
    extracted_text = ""
    rubric_text = ""
    extract_dir = os.path.join(temp_dir, f"extracted_{uuid.uuid4()}")

    # Keywords that might indicate a rubric file
    rubric_keywords = [
        "rubric", "marking", "grading", "criteria",
        "evaluation", "scoring", "marks", "grade"
    ]

    try:
        if file_type == "zip":
            logger.info(f"Extracting and processing files from ZIP: {file_path}")
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)
        elif file_type == "rar":
            logger.info(f"Extracting and processing files from RAR: {file_path}")
            try:
                with rarfile.RarFile(file_path) as rf:
                    rf.extractall(extract_dir)
            except rarfile.Error as e:
                error_msg = f"RAR extraction failed: {str(e)}. Please ensure 'unrar' tool is installed and accessible."
                logger.error(error_msg)
                return "", error_msg
        elif file_type in ["tar", "tar.gz", "tar.bz2"]:
            logger.info(f"Extracting and processing files from TAR: {file_path}")
            with tarfile.open(file_path) as tf:
                tf.extractall(extract_dir)
        else:
            logger.error(f"Unsupported archive type: {file_type}")

        # Process the extracted files
        for root, _, files in os.walk(extract_dir):
            for filename in files:
                ext = filename.split(".")[-1].lower()
                full_path = os.path.join(root, filename)
                logger.info(f"Processing file in archive: {filename}")

                try:
                    text = ""
                    if ext in ["doc", "docx"]:
                        text = extract_text_from_word(full_path)
                        if not text.strip():
                            try:
                                images = word_to_images(full_path, temp_dir)
                            except:
                                images = WORD_TO_IMAGES(full_path, temp_dir)
                            for img in images:
                                text += extract_text_from_image(img) + "\n"
                                os.remove(img)
                    elif ext in ["ppt", "pptx"]:
                        text = extract_text_from_ppt_with_images(full_path)
                    elif ext == "pdf":
                        text = extract_text_from_pdf(full_path)
                        if not text.strip():
                            images = pdf_to_images(full_path, temp_dir)
                            for img in images:
                                text += extract_text_from_image(img) + "\n"
                                os.remove(img)
                    elif ext in ["jpeg", "jfif", "png", "jpg"]:
                        text = extract_text_from_image(full_path)
                    else:
                        logger.warning(f"Unsupported file type inside archive: {filename}")
                        error_message = f"Unsupported file type inside archive: {filename}"
                        file_extension = filename.split(".")[-1].lower() if "." in filename else filename
                        continue

                    # Check if the filename contains rubric-related keywords
                    filename_lower = filename.lower()
                    if any(keyword in filename_lower for keyword in rubric_keywords):
                        rubric_text += f"\n\n--- File: {filename} ---\n{text}"
                    else:
                        extracted_text += f"\n\n--- File: {filename} ---\n{text}"

                except Exception as e:
                    logger.error(f"Error processing file {filename} inside archive: {e}")

    except Exception as e:
        logger.error(f"Error opening or extracting archive file: {e}")
        return "", JSONResponse(status_code=500, content={"error": f"Failed to process archive file: {str(e)}"})

    finally:
        # Cleanup extracted directory
        if os.path.exists(extract_dir):
            shutil.rmtree(extract_dir, ignore_errors=True)

    return extracted_text, rubric_text

def detect_rubric(text: str) -> dict:
    """
    Detects and extracts rubrics/marking schemes from the provided text,
    then returns them as a structured Python dict (parsed from JSON).

    If no rubric is found, returns {"message": "No rubric or marking scheme detected"}.
    In case of an API error, returns {"error": "Error detecting rubric"}.
    """
    prompt = """
    ROLE: RUBRIC DETECTION SPECIALIST
    Your task is to ONLY detect and extract rubrics, marking schemes, and grading criteria from the provided text.

    Focus STRICTLY on:
    1. Grading criteria
    2. Marking schemes
    3. Assessment rubrics
    4. Evaluation criteria
    5. Scoring guidelines
    6. Performance indicators
    7. Mark allocation
    8. Grade boundaries
    9. Assessment breakdowns
    10. Marking distribution

    RULES:
    - ONLY extract information related to grading, marking, or assessment criteria.
    - IGNORE all other content (instructions, requirements, etc.).
    - If no rubric is found, return a JSON object with: {{ "message": "No rubric or marking scheme detected" }}.
    - Format the output as a single JSON object, with the following top-level keys (use exactly these keys, even if some arrays end up empty):
    {{
        "main_criteria": [
        {{
            "criterion_title": "String",
            "description": "String (optional)",
            "sub_criteria": [
            {{
                "title": "String",
                "description": "String (optional)",
                "mark_allocation": "String or number (e.g., '10 pts')"
            }}
            // …repeat for each sub-criterion…
            ],
            "mark_allocation": "String or number (e.g., '20 pts')"
        }}
        // …repeat for each main criterion…
        ],
        "grade_boundaries": [
        {{
            "grade": "String (e.g., 'A', '80–89%')",
            "threshold": "String or number (e.g., '80%')"
        }}
        // …repeat for each boundary…
        ],
        "performance_levels": [
        {{
            "level_name": "String (e.g., 'Excellent', 'Satisfactory')",
            "description": "String"
        }}
        // …repeat for each level…
        ]
    }}

    Text to analyze:
    {text}
    """

    try:
        # Call OpenAI ChatCompletion
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a specialized rubric detection assistant that ONLY extracts grading and marking information."
                },
                {
                    "role": "user",
                    "content": prompt.format(text=text)
                },
            ],
            response_format={"type": "json_object"},  
            temperature=0
        )

        rubric_json_str = response.choices[0].message.content.strip()

        # Attempt to parse the JSON string into a Python dict
        try:
            rubric_data = json.loads(rubric_json_str)
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON from model output: {e}\nOutput was:\n{rubric_json_str}")
            # Return a minimal fallback indicating parsing failed
            return {
                "error": "Model output was not valid JSON",
                "raw_output": rubric_json_str
            }

        return rubric_data

    except Exception as e:
        logger.error(f"Error in rubric detection: {str(e)}", exc_info=True)
        return {"error": "Error detecting rubric"}

def extract_information_with_openai(text: str) -> dict:
    """
    Sends the extracted text to OpenAI's gpt-4o-mini to extract specific details.
    Returns a dictionary with the extracted information.
    """
    # First detect the rubric
    rubric_content = detect_rubric(text)
    
    prompt = f"""
    Extract the following details from the text below. For any information that is not found, leave the response as an empty message.
    1. **Paper Topic/Subject/Module**
    2. **Assignment Type**
        Identify and classify the type of assignment from the following categories:
        - Essay
        - Business Report
        - Presentation Deck
        - Research Proposal
        - Thesis Chapter
        - Poster
        - Lab Report
        - Case Study
        - Literature Review
        - Technical Documentation
        - Other (specify)
        
        If multiple types are mentioned, select the most prominent one.
        If no clear type is mentioned, classify as "Other" with a brief description.

    3. **Deadline**
        - Extract and identify any specific deadlines, time windows, or time-related conditions mentioned in the provided text. Focus on highlighting critical phrases such as late submission windows, estimated completion times, extension policies, or any other time-related details that might imply a deadline.
        Instructions:
        - Explicit Deadline Identification:
        - Look for direct mentions of deadlines, due dates, or submission dates.
        - If a specific date and time are provided, present them in ISO 8601 format (UTC), e.g., 2025-04-24T10:29:50.020Z.
        - If only a date is mentioned without a specific time, format it as YYYY-MM-DDT00:00:00.000Z.

        Relative or Vague Deadlines:
        - If the deadline is expressed relatively or vaguely (e.g., "Friday Week 12 11:30pm", "Due next Friday", "Submission window of 7 calendar days") or is not an explicit calendar date/time, do NOT infer or convert it to a fixed calendar date or ISO format.
        - Instead, output the deadline exactly as written in the source text as a string.

        Contextual Deadline Indicators:
        - Identify phrases that suggest timeframes or durations, such as:
        - "Approximately 5 hours" (indicates expected time commitment)
        - "Submission window of 7 calendar days"
        - Relative time references like "Due next Friday"
        - Mentions of "deadline", "due date", "submit by", "hand in by"
        - Terms like "week 5", "before the lecture", "after the tutorial"

        Inference:
        - Only perform inference or conversion to ISO format if there is clear explicit calendar date/time information.
        - Otherwise, preserve the original text for relative deadlines.

        Output Formatting:
        - Present each identified deadline or time-related condition on a new line.
        - Use bullet points for clarity.
        - Ensure all dates are in ISO 8601 format (UTC) as specified above.
        - Do not include example formats in the output.
        Note:
        - Parse time expressions including 12-hour clock with am/pm, optionally preceded by symbols such as '@' or 'at'.
        - Recognize and parse times given in 12-hour format with am/pm notation (e.g., "5pm", "11:59 pm", "8:30 am").  
        - Convert parsed times accurately to 24-hour format and then to UTC.
        - If no deadlines or time-related conditions are identified in the text, do not include any examples or placeholders in the output. 
        - If a specific date and time include a time zone abbreviation (e.g., BST, EST, PST), **convert the time to UTC by adjusting the time accordingly before formatting as ISO 8601 UTC.**  For example, "18:00 BST" should be converted to "17:00 UTC" before output.

    4. **Word Count** (you have to give only wordcount in numbers like 1000 don't make any other thing like 900-1000 words)
    5. **Number of Pages**
        - Only report page count if the assignment explicitly requires or limits pages.
        - Ignore incidental or footer/header page numbers like Page 1 of 6, Page 2 of 6 and so on.
    6. **Parts in Assignment**
        Look for sections that might be formatted with:
            - Bullet points or numbered lists
            - Sections marked with symbols like "■" or "▪"
            - Tasks prefixed with "Download", "Import", "Complete", "Build", "Summarize", etc.
            - Any steps or tasks such as Task 1, Task 2, Task 3, etc. listed in sequence
    7. **Software Required**
        Look for software names and task summary in 1 or 2 words such as:
            - R, Python (ETL using PonyORM)
            - ASP.NET Core, EntityFramework Core, MSSQL, Visual Studio
            - WITNESS Simulation Model, SAS Visual Analytics
            - XAMPP, Web Servers, SQL Databases
            - HTML, CSS, JavaScript, AJAX, PHP, MySQL
            - Python (Jupyter Notebook, Pandas, Bokeh, SQLAlchemy)
            - Adobe XD, JavaFX, Java, MongoDB
            - Azure services, Cisco Packet Tracer, Wireshark
            - Microsoft Visio, Draw.io, CAD software
            - MATLAB, Fusion 360, Ansys
            - Any other mentioned software like UML Diagrams etc
        Note:
        - Dont consider these techniques as Software Required
            "Differential Scanning Calorimetry (DSC)",
            "X-Ray Diffraction (XRD)",
            "Fourier Transform Infrared Spectroscopy (FTIR)",
            "Scanning Electron Microscopy (SEM)",
            "High-Performance Liquid Chromatography (HPLC)"
         
    8. **Software Version Required**
    9. **Programming Language/Framework Versions**
    10. **Data File Format**
    11. **Code Execution Environment**
    12. **Expected Output Format**
        Also Look for mentions of:
            - Word or PDF document.
            - Assessment Checklist
            - Report
            - Source Code
            - Files (media)
    13. **Integration with Other Tools**
    14. **Dataset Mentioned**
        Look for mentions of:
            - Canvas/Portal provided data
            - Blackboard datasets
            - Simulation model files
            - Any other specified datasets
    15. **University Name**
        - Extract the name of the university from the provided text.
        - If multiple universities are mentioned, select the most relevant one.
        - If no clear university is mentioned, classify as "Other" with a brief description.
    16. **Detailed Instructions**
            ## ROLE: COMPREHENSIVE INSTRUCTIONS EXTRACTOR
            Extract ALL instructions, rubrics, learning outcomes, and assessment details from the text.
            Combine them into a comprehensive set of instructions.
            
            Include:
            - Step-by-step instructions
            - Submission procedures
            - Formatting requirements
            - Technical specifications
            - Process guidelines
            - Important notes
            - Special requirements
            - Procedural details
            - Assessment guidelines
            - Evaluation standards
            - Performance indicators
            - Quality expectations
            - Success criteria
            - Learning outcomes
            - Module objectives
            - Skill development goals
            - Knowledge acquisition targets
            - Competency requirements
            
            Format the output as a comprehensive, well-structured set of instructions.
            Use clear headings and subheadings to organize different types of information.
            Maintain all specific details, numbers, and requirements exactly as stated.
            
            Return as a summarized information.
    17. Citation style
        - APA
        - MLA
        - Chicago
        - Harvard
        - Vancouver
        - IEEE
        NOTE Also return the total number of citations required along citation style 
            Look for phrases like:
            - "between 6 and 10 references"
            - "minimum 5 references"
            - "at least 8 sources"
            - "6-10 citations"
            - "references should be between X and Y"
            Extract both the citation style and the total number of references required. 
    **Text:**
    {text}

    **Respond only with the extracted details in JSON format using the following keys:**
    - paper_topic
    - assignment_type
    - deadline
    - word_count
    - page_count
    - assignment_parts
    - software_required
    - software_version
    - programming_language_version
    - data_file_format
    - execution_environment
    - expected_output
    - integration_tools
    - dataset_mentioned
    - university_name
    - instructions 
    - citation_style
    - total_references_required
    """

    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are an assistant that extracts structured information."},
                {"role": "user", "content": prompt},
            ],
            response_format={"type": "json_object"},
            temperature=0
        )
        response_content = response.choices[0].message.content
        
        # Parse the response as JSON
        extracted_data = json.loads(response_content)
        
        # Combine the rubric content with the instructions
        instructions = extracted_data.pop("instructions", "")
        combined_instructions = f"DETAILED INSTRUCTIONS:\n{instructions}"
        extracted_data["_instructions"] = combined_instructions
        extracted_data["_rubric_json"] = rubric_content
        
        return extracted_data

    except json.JSONDecodeError:
        logger.error("Failed to parse OpenAI response as JSON.", exc_info=True)
        return {"error": "Failed to parse extracted data."}
    except Exception as e:
        logger.error(f"OpenAI API error: {str(e)}", exc_info=True)
        return {"error": str(e)}

def get_static_instructions(assignment_type: str) -> str:
    """
    Returns static instructions based on assignment type.
    """
    static_instructions = {
        "Essay": """Essay Structure:
            • Introduction
            •	Clearly define the topic scope (what is included/excluded).
            •	Explain why the topic matters—academic, social, professional relevance.
            •	Introduce the thesis statement—your central argument in 1–2 concise sentences.
            •	Outline key themes or sections in the order they appear.
            • Main Body
            •	Each paragraph = one clear argument or point tied to the thesis.
            •	Begin with a topic sentence that signals the paragraph's main idea.
            •	Include evidence: research, examples, stats, or theory.
            •	Add critical analysis: compare, contrast, interpret sources.
            •	Use signposting and transitions.
            • Conclusion
            •	Rephrase the thesis.
            •	Summarise key arguments made.
            •	Reflect on the broader implications.
            • Other Notes
            •	Usually no headings unless specified.
            •	Use third-person academic tone.
            •	Maintain spelling consistency and correct academic tense logic.""",

        "Business Report": """Business Report Structure:
            • Executive Summary
            •	Write last but placed first—summarises the whole report in 150–250 words.
            •	Include: context, objectives, methods, key findings, and recommendations.
            • Introduction
            •	Define the problem or opportunity.
            •	Clearly state objectives or aims.
            •	Provide a brief background.
            • Main Body
            •	Organise using section headings.
            •	Include data and evidence.
            •	Link findings to models/frameworks.
            • Recommendations
            •	Actionable, bullet-style suggestions.
            •	Justify each recommendation.
            • Conclusion
            •	Briefly restate findings.
            •	Emphasise the business value.""",

        "Report": """Business Report Structure:
            • Executive Summary
            •	Write last but placed first—summarises the whole report in 150–250 words.
            •	Include: context, objectives, methods, key findings, and recommendations.
            • Introduction
            •	Define the problem or opportunity.
            •	Clearly state objectives or aims.
            •	Provide a brief background.
            • Main Body
            •	Organise using section headings.
            •	Include data and evidence.
            •	Link findings to models/frameworks.
            • Recommendations
            •	Actionable, bullet-style suggestions.
            •	Justify each recommendation.
            • Conclusion
            •	Briefly restate findings.
            •	Emphasise the business value.""",

        "Research Proposal": """Research Proposal Structure:
            • Introduction
            •	Define the research problem or gap.
            •	Explain its importance.
            •	Provide clear research question(s).
            • Literature Review
            •	Organise sources by theme or chronology.
            •	Identify areas of agreement/disagreement.
            • Methodology
            •	Specify research design.
            •	Describe sampling, tools, methods.
            •	Include data analysis plans.
            • Timeline & Budget
            •	Provide realistic time allocation.
            •	Include budget line items.""",

        "Presentation Slides": """Presentation Structure:
            • Introduction Slide
            •	Include title, name, date, affiliation.
            •	Add purpose statement.
            • Agenda Slide
            •	List presentation structure.
            • Main Content Slides
            •	One core idea per slide.
            •	Use concise bullet points.
            •	Include relevant visuals.
            • Summary Slide
            •	Reiterate key points.
            •	Emphasise conclusions.""",

        "Poster": """Poster Structure:
            • Introduction
            •	Articulate research question.
            •	Use bold text or icons.
            • Methods
            •	Describe methodology concisely.
            • Results
            •	Use visuals over words.
            • Discussion
            •	Interpret findings briefly.
            • Conclusion
            •	Highlight main insight.""",

        "Reflective Writing": """Reflective Writing Structure:
            • Introduction
            •	Describe the event/experience.
            •	Explain its importance.
            • Body
            •	Use a structured model.
            •	Break down into description, feelings, evaluation.
            •	Link experiences to theory.
            • Conclusion
            •	State future applications.
            •	Reflect on growth.""",

        "Annotated Bibliography": """Annotated Bibliography Structure:
            • Citation Entry
            •	Use correct citation format.
            •	Alphabetically arranged.
            • Summary
            •	Describe main idea, method, findings.
            • Evaluation
            •	Comment on relevance, credibility.""",

        "Literature Review": """Literature Review Structure:
            • Introduction
            •	Define research topic.
            •	Summarise search strategy.
            • Main Body
            •	Group literature by theme.
            •	Compare and critique sources.
            •	Highlight consensus and gaps.
            • Conclusion
            •	Summarise key insights.
            •	Show research direction."""
    }
    
    return static_instructions.get(assignment_type, "No specific instructions available for this assignment type.")

# Add this new function before the existing langchain function:
def create_enhanced_section_connection_prompt(current_section, previous_sections, module_materials, assignment_brief):
    """Create prompts that ensure sections are interconnected"""
    
    connection_prompt = f"""
SECTION INTERCONNECTION STRATEGY:
- Current Section: {current_section}
- Previous Sections Covered: {previous_sections}
- Connect this section to previous arguments and findings
- Reference back to earlier points where relevant
- Build upon established concepts from module materials
- Maintain logical flow and narrative consistency

MODULE MATERIALS INTEGRATION:
{module_materials}

ASSIGNMENT BRIEF CONTEXT:
{assignment_brief}

Ensure this section:
1. References relevant concepts from module materials
2. Connects to previous section arguments
3. Builds logical progression toward assignment objectives
4. Uses appropriate transitions and linking phrases
"""
    return connection_prompt

def generate_assignment_outline_with_langchain(extracted_text: str, module_materials: str, additional_information: str) -> str:
    """
    Generates an assignment outline using the 7-step LangChain workflow.
    """

    llm = ChatOpenAI(temperature=0, model="gpt-4o")
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)

    # Step 1: Assignment Brief Clarity
    prompt1 = PromptTemplate(
        input_variables=["assignment_context", "chat_history"],
        template="""
{chat_history}
Prompt 1: Assignment Brief Clarity
Act as an academic assignment interpreter with expertise in university brief analysis. Based on the uploaded assignment brief and any associated module materials, reproduce the full task in detail.
Your response should clearly outline the following:
• Assignment title
• Module name and code (if available)
• Type of assessment (e.g., essay, report, case study)
• Exact task instructions (use full phrasing from brief)
• Word count and formatting requirements
• Submission deadline
• Referencing style required (e.g., APA, Harvard)
If multiple documents are uploaded, extract and integrate all relevant details into a single coherent brief summary. Avoid interpretation — quote original phrasing where possible.

Assignment Context:
{assignment_context}
"""
    )
    chain1 = LLMChain(llm=llm, prompt=prompt1, memory=memory, output_key="clarified_brief")

    # Step 2: Deliverables Breakdown
    prompt2 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
{chat_history}
Prompt 2: Deliverables Breakdown
Act as an academic assignment planner trained in deliverable decomposition and assessment structuring. From the assignment brief, identify and list all expected deliverables. Clarify what specific outputs the writer is required to produce.
Your response should include:
• The primary deliverable (e.g., a 2,000-word essay/Report etc)
• Any embedded components (e.g., executive summary, literature review, SWOT analysis, reflection section)
• Formatting or file-type requirements (e.g., Word doc, pdf, PPT etc)
• Special instructions (e.g., inclusion of tables, models, appendices, or specific sections)
If the brief includes multiple tasks or parts (e.g., Part A and Part B), list them separately and describe what each one entails.
"""
    )
    chain2 = LLMChain(llm=llm, prompt=prompt2, memory=memory, output_key="deliverables")

    # Step 3: Module Files Strategy
    prompt3 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
        {chat_history}
        Prompt 3: Module Files Strategy
        Act as a module-savvy academic writing advisor with expertise in applying lecture content and course theories to academic assignments. Review all uploaded module materials — such as lecture slides, topic summaries, or handbooks — and extract any relevant frameworks, models, or theories that could support this assignment. Based on the assignment topic (e.g., globalization), identify which concepts from the module should be incorporated and where.
        Your response should include:
        • A list of module files reviewed (e.g., "Lecture 5: Globalization Frameworks")
        • Theories or concepts mentioned in them (e.g., Ghemawat's CAGE framework, Global Value Chain theory)
        • For each concept, explain:
            o What it means
            o How and where I could use it in my assignment (e.g., to structure a section, support analysis, or offer critique)
        If the module content doesn't directly mention applicable theories, state that clearly and suggest a general academic approach (e.g., independent research, external frameworks). Prioritize content that aligns with the learning outcomes or rubric.
        """
    )
    chain3 = LLMChain(llm=llm, prompt=prompt3, memory=memory, output_key="module_strategy")

    # Step 4: Learning Outcomes Focus
    prompt4 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
        {chat_history}
        Prompt 4: Learning Outcomes Focus
        Act as a learning outcome alignment strategist with expertise in helping academic writers meet university assessment goals. Identify and interpret the learning outcomes listed in the assignment brief or module guide. Explain what each outcome requires me to demonstrate, and provide specific guidance on how to reflect those skills or competencies in my assignment.
        Your response should include:
            • A list of all relevant learning outcomes (quoted directly from the brief if available)
            • For each one:
            o A plain-language explanation of what it means (e.g., "evaluate theories" = critically compare and judge theories)
            o Suggestions for where in the assignment to demonstrate that skill (e.g., "Outcome 2 should be addressed in your analysis section using X framework")
        If no explicit learning outcomes are given, infer possible expectations based on the assignment type, topic, and rubric — and provide guidance accordingly.
        """
    )
    chain4 = LLMChain(llm=llm, prompt=prompt4, memory=memory, output_key="learning_outcomes")

    # Step 5: Strategic Use of Module Reading List
    prompt5 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
        {chat_history}
        Prompt 5: Strategic Use of Module Reading List
        Act as a university-level academic research assistant with a focus on integrating module-specific readings into assignment writing with precision and relevance.
        Analyze all uploaded materials, including weekly resources, topic folders, and separate reading list files. From these, extract only the readings that are directly applicable to the assignment topic (e.g., globalization, strategic management, etc.). Do not fabricate or infer sources — only include readings that are explicitly listed by the module instructor.
        Your response must include:
            • Exact titles and author names of the readings as listed in the module materials (e.g., "Johnson et al., Exploring Strategy, Week 5 reading")
            • Type of source (e.g., journal article, book chapter, report)
            • A brief explanation of the reading's core relevance to the assignment topic
            • Direct suggestions for how each reading can be cited in the assignment (e.g., "use this source to support your critique of globalization models in the second section")
        If the module provides multiple readings, prioritize those that are most clearly aligned with the assignment task. If no relevant readings are available or identifiable, clearly state that and advise whether external peer-reviewed sources may be necessary to support key arguments.
        """
    )
    chain5 = LLMChain(llm=llm, prompt=prompt5, memory=memory, output_key="reading_list")

    # Step 6: Marking Rubric Translator & Grade Maximizer
    prompt6 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
        {chat_history}
        Prompt 6: Marking Rubric Translator & Grade Maximizer
        Act as a university assessment expert and academic writing strategist with deep knowledge of grading rubrics and high-mark criteria. Review the marking rubric or grading criteria provided in the assignment brief or module guide. Translate each assessment criterion into practical writing advice and show how to structure the assignment to meet the top-level performance descriptors (e.g., "Excellent", "Distinction").
        Your response must include:
            • A clear summary of each rubric category (e.g., Argument Quality, Use of Evidence, Structure and Flow, Critical Analysis, Referencing, Academic Style)
            • An explanation of what the highest grade band looks like for each category
            • Specific instructions on how the writer should construct the assignment to meet or exceed those standards
                o (e.g., "To achieve top marks in Structure, organize your report using clear section headings aligned to the task and rubric expectations.")
        If the actual rubric is missing, infer likely criteria based on UK academic standards and label it accordingly. Do not fabricate any rubric text — only summarize what's visible or standard. Provide actionable, performance-focused writing advice that directly supports a high-grade outcome.
        """
    )
    chain6 = LLMChain(llm=llm, prompt=prompt6, memory=memory, output_key="rubric_advice")

    # Step 7: Strategic Outline Builder (Writer-Directed)
    prompt7 = PromptTemplate(
        input_variables=["input", "chat_history"],
        template="""
        {chat_history}
        Prompt 7: Strategic Outline Builder (Writer-Directed)
        Act as a university-level assignment planner and academic structure expert, specialized in rubric-aligned and outcome-driven writing frameworks. Based on everything we've discussed so far — including the assignment brief, expected deliverables, module content, learning outcomes, key readings, and the full marking rubric — generate a comprehensive and detailed outline for this assignment.
        Your outline must serve as a complete blueprint for high-quality academic writing. Please include:
            • Proper section headings tailored to the assignment type (e.g., essay, report, case study)
            • Suggested word count distribution for each section
            • What each section should cover, referencing:
                o Theories or models from the module
                o Specific readings identified earlier
                o Relevant learning outcomes
                o Rubric criteria (e.g., critical analysis, originality, structure)
        Make this outline so detailed and targeted that I can follow it section-by-section to write an assignment that meets all academic expectations and achieves the highest grade band. Do not generalize — incorporate all prior information from this chat thread directly into the structure.
        NOTE: Do not include any other text in your response except the outline.
        """
    )
    chain7 = LLMChain(llm=llm, prompt=prompt7, memory=memory, output_key="final_outline")

    # --- RUN THE FULL WORKFLOW ---
    try:
        assignment_context = (
            f"Assignment Brief:\n{extracted_text}\n\n"
            f"Module Materials:\n{module_materials}\n\n"
            f"Additional Information:\n{additional_information}"
        )
        result1 = chain1.invoke({"assignment_context": assignment_context})
        result2 = chain2.invoke({"input": ""})
        result3 = chain3.invoke({"input": ""})
        result4 = chain4.invoke({"input": ""})
        result5 = chain5.invoke({"input": ""})
        result6 = chain6.invoke({"input": ""})
        result7 = chain7.invoke({"input": ""})
        return result7["final_outline"]
    except Exception as e:
        import traceback
        logger.error(f"Error in LangChain outline workflow: {str(e)}\n{traceback.format_exc()}")
        return "Error generating assignment outline."

def process_helping_material(zip_path: str, request_temp_dir: str) -> str:
    """
    Process the helping material zip file and extract its contents.
    Returns a structured string containing all the helping material content.
    """
    extracted_content = ""
    extract_dir = os.path.join(request_temp_dir, f"helping_material_{uuid.uuid4()}")

    try:
        # Create extraction directory
        os.makedirs(extract_dir, exist_ok=True)

        # Extract the zip file
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)

        # Process each file in the extracted directory
        for root, _, files in os.walk(extract_dir):
            for filename in files:
                file_path = os.path.join(root, filename)
                file_ext = filename.split('.')[-1].lower()

                try:
                    file_content = ""
                    if file_ext in ['doc', 'docx']:
                        file_content = extract_text_from_word(file_path)
                        if not file_content.strip():
                            try:
                                images = word_to_images(file_path, request_temp_dir)
                            except:
                                images = WORD_TO_IMAGES(file_path, request_temp_dir)
                            for img in images:
                                file_content += extract_text_from_image(img) + "\n"
                                os.remove(img)
                    elif file_ext in ['pdf']:
                        file_content = extract_text_from_pdf(file_path)
                    elif file_ext in ['ppt', 'pptx']:
                        file_content = extract_text_from_ppt_with_images(file_path)
                    elif file_ext in ['jpg', 'jpeg', 'png', 'jfif']:
                        file_content = extract_text_from_image(file_path)
                    else:
                        logger.warning(f"Unsupported file type in helping material: {filename}")
                        continue

                    if file_content.strip():
                        extracted_content += f"\n\n--- File: {filename} ---\n{file_content}"

                except Exception as e:
                    logger.error(f"Error processing file {filename} in helping material: {str(e)}")
                    continue

    except Exception as e:
        logger.error(f"Error processing helping material zip: {str(e)}")
        return ""

    finally:
        # Cleanup
        try:
            if os.path.exists(extract_dir):
                shutil.rmtree(extract_dir)
            if os.path.exists(zip_path):
                os.remove(zip_path)
        except Exception as e:
            logger.error(f"Error cleaning up helping material files: {str(e)}")

    return extracted_content

def summarize_text_adaptive(text: str) -> str:
    """
    Summarizes the given text using OpenAI, with summary length based on input size.
    """
    if not text.strip():
        return ""
    word_count = len(text.split())
    # Adaptive summary length
    if word_count < 500:
        summary_instruction = "Summarize the following helping material, keeping all important details. Be concise but do not omit key points."
    elif word_count < 2000:
        summary_instruction = "Summarize the following helping material, focusing on the most important instructions, requirements, and guidelines. Omit repetitive or generic information."
    else:
        summary_instruction = "Summarize the following helping material, condensing it to only the most essential instructions, requirements, and guidelines. Omit all non-essential details and repetition."
    prompt = f"{summary_instruction}\n\nText:\n{text}"
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",  # Use a fast model for summarization
            messages=[
                {"role": "system", "content": "You are a helpful assistant that summarizes academic materials."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1200,  # Allow enough space for longer summaries if needed
            temperature=0.3
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error summarizing helping material: {str(e)}", exc_info=True)
        return text  # Fallback to original if summarization fails

# Add a function to determine file type from extension
def get_file_type_from_extension(filename: str) -> str:
    """Determine file type from file extension."""
    if not filename:
        return ""
    # Get the extension and remove the dot
    ext = filename.split('.')[-1].lower()
    # Handle special cases for tar files
    if ext in ['gz', 'bz2'] and filename.lower().endswith('.tar.' + ext):
        return f"tar.{ext}"
    return ext

class ExtractRequest(BaseModel):
    customer_id: str

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],        # <- allow any origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/extract/")
async def extract_data(
    request: Request,
    file: UploadFile = File(None),  # file upload is optional
    helping_material: UploadFile = File(None),  # New parameter for helping material zip
    additional_information: str = Form(None),  # New parameter for additional information
):
    # Create user session
    session = create_user_session()
    session_id = session.session_id
    request_temp_dir = session.temp_dir
    file_path = ""
    helping_material_path = ""
    helping_material_content = ""
    
    try:
        # Process helping material if provided
        if helping_material:
            logger.info(f"Processing helping material zip: {helping_material.filename}")
            helping_material_path = os.path.join(request_temp_dir, f"helping_material_{uuid.uuid4()}.zip")
            with open(helping_material_path, "wb") as buffer:
                shutil.copyfileobj(helping_material.file, buffer)
            helping_material_content = process_helping_material(helping_material_path, request_temp_dir)
            logger.info("Helping material processed successfully")

        # Handle file upload
        if file:
            logger.info(f"Processing uploaded file: {file.filename}")
            # Determine file type from extension
            file_type = get_file_type_from_extension(file.filename)
            
            # Validate file type
            if file_type not in ALLOWED_FILE_TYPES:
                logger.warning(f"Unsupported file type: {file_type}")
                error_message = f"Unsupported file type '{file_type}'. Allowed types are: {', '.join(ALLOWED_FILE_TYPES)}."
                return JSONResponse(
                    status_code=400, 
                    content={"error": error_message}
                )
            
            file_path = os.path.join(request_temp_dir, f"temp_{uuid.uuid4()}.{file_type}")
            original_filename = file.filename
            with open(file_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            logger.info(f"File uploaded to: {file_path}")
        else:
            return JSONResponse(status_code=400, content={"error": "No file provided."})

        # Determine actual_filename and document_path
        if file:
            actual_filename = file.filename
            document_path = file_path  # This is the temp path used for processing
        else:
            actual_filename = ""
            document_path = ""

        # Extract text based on file type
        if file_type in ["ppt", "pptx"]:
            logger.info(f"Extracting text from PPT/PPTX: {file_path}")
            extracted_text = extract_text_from_ppt_with_images(file_path)
        elif file_type in ["doc", "docx"]:
                try:
                    logger.info(f"Attempting to extract text from Word document: {file_path}")
                    extracted_text = extract_text_from_word(file_path)
                except Exception as e:
                    logger.error(f"Error extracting text from Word document: {e}", exc_info=True)

                if not extracted_text.strip():
                    logger.info("Initial text extraction failed or returned empty. Attempting OCR on images.")

                    try:
                        image_paths = word_to_images(file_path, request_temp_dir)
                    except Exception as e:
                        logger.error(f"Error converting Word document to images: {e}", exc_info=True)
                        image_paths = []

                    if not image_paths:
                        logger.info("word_to_images failed. Attempting WORD_TO_IMAGES as a fallback.")
                        try:
                            image_paths = WORD_TO_IMAGES(file_path, request_temp_dir)
                        except Exception as e:
                            logger.error(f"Error in WORD_TO_IMAGES fallback: {e}", exc_info=True)
                            image_paths = []

                    for img_path in image_paths:
                        try:
                            ocr_text = extract_text_from_image(img_path)
                            extracted_text += ocr_text + "\n"
                        except Exception as e:
                            logger.error(f"Error extracting text from image {img_path}: {e}", exc_info=True)
                        finally:
                            try:
                                os.remove(img_path)
                            except Exception as e:
                                logger.warning(f"Could not delete image {img_path}: {e}")

                    try:
                        shutil.rmtree("output_images", ignore_errors=True)
                    except Exception:
                        pass
                    
        elif file_type == "pdf":
            logger.info(f"Extracting text from PDF: {file_path}")
            extracted_text = extract_text_from_pdf(file_path)
            # If extraction returns no text, fall back to image conversion and OCR
            if not extracted_text.strip():
                logger.info("PDF text extraction returned empty. Converting PDF pages to images for OCR.")
                image_paths = pdf_to_images(file_path, request_temp_dir)  # Convert PDF pages to images
                extracted_text = ""
                for img_path in image_paths:
                    ocr_text = extract_text_from_image(img_path)  # Extract text via OCR
                    extracted_text += ocr_text + "\n"
                    os.remove(img_path)  # Remove the individual image
                # Optionally, remove the output folder if it's empty
                try:
                    os.rmdir("output_images")
                except Exception:
                    pass

        elif file_type in ["jpeg", "jfif", "png", "jpg"]:
            logger.info(f"Extracting text from image: {file_path}")
            extracted_text = extract_text_from_image(file_path)

        elif file_type in ["zip", "rar", "tar", "tar.gz", "tar.bz2"]:
            extracted_text, rubric_text = extract_from_archive(file_path, file_type, request_temp_dir)
            
            # If we found a separate rubric file, use its content
            if rubric_text:
                rubric_content = detect_rubric(rubric_text)
            else:
                # Fall back to detecting rubric in the main text if not found separately
                rubric_content = {"message": "No rubric or marking scheme detected"}

        else:
            logger.error(f"Unsupported file type encountered: {file_type}")
            error_message = f"Unsupported file type: {file_type}"
            raise ValueError(error_message)
        
        if not isinstance(extracted_text, str):
            logger.error("Extracted text is not a string.")
            error_message = "Extracted text is not a string"
            return JSONResponse(status_code=500, content={"error": "Internal server error."})

        if not extracted_text.strip():
            logger.warning("No text extracted from the file.")
            error_message = "No text could be extracted from the file"
            return JSONResponse(status_code=400, content={"error": "No text could be extracted from the file."})

        logger.debug(f"Extracted Text (Brief): {extracted_text[:500]}...")
        
        # Process the extracted text with OpenAI for detailed extraction
        # This should run for all file types, not just archives
        processed_data = extract_information_with_openai(extracted_text)
        
        # For archive files, we already have rubric_content, for others we need to get it from processed_data
        if file_type not in ["zip", "rar", "tar", "tar.gz", "tar.bz2"]:
            rubric_content = processed_data.get("_rubric_json", {"message": "No rubric or marking scheme detected"})

        # Use LangChain workflow to generate the outline instead of draft
        assignment_outline = generate_assignment_outline_with_langchain(
            extracted_text, helping_material_content, additional_information or ""
        )

        # Store assignment data in database
        try:
            # First insert the basic assignment record
            assignment_id = db.insert_assignment(
                assignment_type=processed_data.get("assignment_type", "Other")
            )

            # Then insert the detailed information
            def safe_str(val):
                if isinstance(val, list):
                    return ", ".join(str(v) for v in val)
                return str(val) if val is not None else ""

            detail_id = db.insert_assignment_detail(
                assignment_id=assignment_id,
                word_count=int(processed_data.get("word_count", "0") or "0"),
                due_date=safe_str(processed_data.get("deadline", "")),
                assignment_type=safe_str(processed_data.get("assignment_type", "Other")),
                software_required=safe_str(processed_data.get("software_required", "")),
                topic=safe_str(processed_data.get("paper_topic", "")),
                university_name=safe_str(processed_data.get("university_name", "")),
                citation_style=safe_str(processed_data.get("citation_style", ""))
            )

            # Insert instructions with additional information
            db.insert_assignment_instruction(
                assignment_id=assignment_id,
                instruction=processed_data.pop("_instructions", ""),
                rubric=json.dumps(rubric_content),
                static_instruction=get_static_instructions(processed_data.get("assignment_type", "Other")),
                additional_information=additional_information if additional_information else ""
            )

            # Insert material
            db.insert_assignment_material(
                assignment_id=assignment_id,
                actual_filename=actual_filename,
                document_path=document_path,
                helping_material=helping_material_content
            )

            # Store the outline instead of draft
            db.insert_assignment_text(
                assignment_id=assignment_id,
                filename=actual_filename,
                brief=extracted_text,
                draft=assignment_outline  # Store as draft for now, or rename column to 'outline' if needed
            )
            processed_data["assignment_id"] = assignment_id    
            processed_data["session_id"] = session_id
            return {
                "data": processed_data, 
                "message": "Assignment data and outline generated and stored.",
                "status": "complete",
                "session_id": session_id
            }

        except mysql.connector.Error as err:
            logger.error(f"Database error: {err}")
            return JSONResponse(status_code=500, content={"error": "Database error occurred"})
            
    except Exception as e:
        error_message = str(e)
        logger.error(f"Error processing file for customer : {error_message}", exc_info=True)
        return JSONResponse(status_code=500, content={"error": str(e)})
    finally:
        # Schedule session cleanup after some time
        threading.Timer(300, cleanup_user_session, args=[session_id]).start()  # 5 minutes

# ENHANCED SUBSECTION-WISE ASSIGNMENT GENERATION ENDPOINT
@app.post("/generate-assignment-subsections/")
async def generate_assignment_subsections(request: AssignmentGenerationRequest):
    """
    Generate complete assignment using SUBSECTION-WISE methods for better word count compliance
    """
    try:
        assignment_id = request.assignment_id
        
        # Retrieve assignment data from database
        assignment_data = db.get_assignment_by_id(assignment_id)
        if not assignment_data:
            return JSONResponse(
                status_code=404, 
                content={"error": f"Assignment with ID {assignment_id} not found"}
            )
        
        # Extract required data from database
        assignment_brief = assignment_data.get("brief", "")
        assignment_outline = assignment_data.get("outline", "")  # This was stored as 'draft'
        helping_material = assignment_data.get("helping_material", "")
        assignment_type = assignment_data.get("assignment_type", "Other")
        
        if not assignment_brief:
            return JSONResponse(
                status_code=400,
                content={"error": "Assignment brief not found in database"}
            )
        
        if not assignment_outline:
            return JSONResponse(
                status_code=400,
                content={"error": "Assignment outline not found in database"}
            )
        
        logger.info(f"🚀 Starting SUBSECTION-WISE assignment generation for ID: {assignment_id}")
        
        # Step 1: Enhanced assignment structure extraction
        logger.info("📋 Step 1: Enhanced assignment structure extraction...")
        json_output = extract_assignment_structure_advanced(assignment_brief + "\n\n" + assignment_outline)
        
        if not json_output:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to extract assignment structure using enhanced method"}
            )
        
        # Step 2: Save structure with enhanced validation
        logger.info("📋 Step 2: Saving assignment structure with enhanced validation...")
        structure_data = save_assignment_structure_advanced(json_output, assignment_id)
        
        if not structure_data:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to save assignment structure"}
            )
        
        # Step 3: Create SUBSECTION-WISE prompts with precise targeting
        logger.info("📋 Step 3: Creating SUBSECTION-WISE prompts with precise targeting...")
        subsection_prompts = create_advanced_subsection_prompts(
            structure_data, 
            assignment_brief, 
            helping_material, 
            assignment_id
        )
        
        if not subsection_prompts:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to create subsection-wise prompts"}
            )
        
        # Step 4: Generate all subsections with precise word count compliance
        logger.info("📋 Step 4: Generating assignment SUBSECTIONS with precise word count compliance...")
        generated_files = generate_all_assignment_subsections_advanced(subsection_prompts, assignment_id)
        
        if not generated_files:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to generate assignment subsections"}
            )
        
        # Step 5: Create professional DOCX from subsections
        logger.info("📄 Step 5: Creating professional DOCX from subsections...")
        docx_file = create_professional_docx_from_subsections(assignment_id, structure_data)
        
        if not docx_file:
            return JSONResponse(
                status_code=500,
                content={"error": "Failed to create professional DOCX document"}
            )
        
        # Calculate enhanced statistics
        estimated_words = 0
        total_paragraphs = 0
        citations_found = 0
        subsections_generated = 0
        sections_generated = 0
        
        if os.path.exists(docx_file):
            try:
                from docx import Document
                doc = Document(docx_file)
                total_paragraphs = len(doc.paragraphs)
                estimated_words = sum(len(p.text.split()) for p in doc.paragraphs if p.text.strip())
                
                # Count citations and subsections
                subsections_dir = os.path.join("generated_assignments", f"assignment_{assignment_id}", "subsections")
                citations = extract_citations_from_subsections(subsections_dir)
                citations_found = len(citations)
                
                # Count subsections vs sections
                for prompt_data in subsection_prompts:
                    if prompt_data.get("is_subsection", False):
                        subsections_generated += 1
                    else:
                        sections_generated += 1
                
            except Exception as e:
                logger.warning(f"Failed to calculate enhanced statistics: {e}")
        
        logger.info(f"🎯 SUBSECTION-WISE assignment generation completed for ID: {assignment_id}")
        logger.info(f"✅ Subsection-wise Features Applied:")
        logger.info(f"   - Generated {subsections_generated} subsections and {sections_generated} sections")
        logger.info(f"   - Each subsection follows exact word count (better compliance)")
        logger.info(f"   - Smart citation control per subsection")
        logger.info(f"   - Enhanced context preservation")
        logger.info(f"   - Professional References section with auto-generated citations")
        logger.info(f"   - Professional formatting throughout")
        
        return {
            "assignment_id": assignment_id,
            "message": "Subsection-wise assignment generated successfully",
            "status": "completed",
            "docx_file": docx_file,
            "subsections_generated": subsections_generated,
            "sections_generated": sections_generated,
            "total_parts_generated": len(generated_files),
            "estimated_words": estimated_words,
            "total_paragraphs": total_paragraphs,
            "citations_found": citations_found,
            "generated_files": generated_files,
            "subsection_features": [
                f"Generated {subsections_generated} subsections with precise word count targeting",
                f"Generated {sections_generated} regular sections",
                "Smart citation control per subsection (no citations in intro/summary/conclusion)",
                "Enhanced context preservation for each subsection",
                "Auto-generated references from citations",
                "Professional formatting and structure",
                "Title page and table of contents",
                "Better word count compliance through subsection-wise generation"
            ]
        }
        
    except Exception as e:
        error_message = str(e)
        logger.error(f"Error in subsection-wise assignment generation: {error_message}", exc_info=True)
        
        # Update database with error status
        try:
            db.update_assignment_generation_status(assignment_id, "failed", None, error_message)
        except Exception as db_err:
            logger.error(f"Failed to update database with error status: {db_err}")
        
        return JSONResponse(status_code=500, content={"error": error_message})

# UPDATE THE EXISTING ADVANCED ENDPOINT TO USE SUBSECTION APPROACH
@app.post("/generate-assignment-advanced/")
async def generate_assignment_advanced_updated(request: AssignmentGenerationRequest):
    """
    Enhanced assignment generation endpoint - now uses subsection-wise approach by default
    """
    logger.info(f"Advanced endpoint called, using subsection-wise generation for assignment {request.assignment_id}")
    return await generate_assignment_subsections(request)

# LEGACY ENDPOINT (keeping for backward compatibility)
@app.post("/generate-assignment/")
async def generate_assignment_legacy(request: AssignmentGenerationRequest):
    """
    Legacy assignment generation endpoint - redirects to enhanced method
    """
    logger.info(f"Legacy endpoint called, redirecting to enhanced generation for assignment {request.assignment_id}")
    return await generate_assignment_subsections(request)

@app.get("/assignment/{assignment_id}")
async def get_assignment(assignment_id: int):
    """
    Get assignment data by ID
    """
    try:
        assignment_data = db.get_assignment_by_id(assignment_id)
        if not assignment_data:
            return JSONResponse(
                status_code=404,
                content={"error": f"Assignment with ID {assignment_id} not found"}
            )
        
        return {
            "assignment_id": assignment_id,
            "data": assignment_data,
            "status": "success"
        }
        
    except Exception as e:
        logger.error(f"Error retrieving assignment {assignment_id}: {str(e)}", exc_info=True)
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.get("/download-assignment/{assignment_id}")
async def download_assignment_file(assignment_id: int):
    """
    Download generated assignment DOCX file
    """
    try:
        # Get assignment data to find the generated file
        assignment_data = db.get_assignment_by_id(assignment_id)
        if not assignment_data:
            return JSONResponse(
                status_code=404,
                content={"error": f"Assignment with ID {assignment_id} not found"}
            )
        
        # Check if file exists in the expected location
        assignment_dir = os.path.join("generated_assignments", f"assignment_{assignment_id}")
        
        # Look for DOCX files in the assignment directory
        docx_files = []
        if os.path.exists(assignment_dir):
            for file in os.listdir(assignment_dir):
                if file.endswith('.docx'):
                    docx_files.append(os.path.join(assignment_dir, file))
        
        if not docx_files:
            return JSONResponse(
                status_code=404,
                content={"error": "Generated assignment file not found. Please generate the assignment first."}
            )
        
        # Use the first (and likely only) DOCX file found
        file_path = docx_files[0]
        
        if not os.path.exists(file_path):
            return JSONResponse(
                status_code=404,
                content={"error": "Assignment file not found on disk"}
            )
        
        # Prepare file for download
        from fastapi.responses import FileResponse
        
        # Generate a user-friendly filename
        assignment_type = assignment_data.get('assignment_type', 'Assignment')
        topic = assignment_data.get('topic', '')
        
        # Clean filename
        clean_type = re.sub(r'[^\w\s-]', '', assignment_type).strip()
        clean_topic = re.sub(r'[^\w\s-]', '', topic).strip() if topic else ''
        
        if clean_topic:
            download_filename = f"{clean_type}_{clean_topic}_ID{assignment_id}.docx"
        else:
            download_filename = f"{clean_type}_Assignment_ID{assignment_id}.docx"
        
        # Clean filename further
        download_filename = re.sub(r'[-\s]+', '_', download_filename)
        
        logger.info(f"Serving download for assignment {assignment_id}: {download_filename}")
        
        return FileResponse(
            path=file_path,
            filename=download_filename,
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
        
    except Exception as e:
        logger.error(f"Error downloading assignment {assignment_id}: {str(e)}", exc_info=True)
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.get("/assignment-history")
async def get_assignment_history():
    """
    Get assignment history
    """
    try:
        history = db.get_assignment_history()
        return {
            "history": history,
            "status": "success"
        }
        
    except Exception as e:
        logger.error(f"Error retrieving assignment history: {str(e)}", exc_info=True)
        return JSONResponse(status_code=500, content={"error": str(e)})

# UTILITY FUNCTIONS FOR ENHANCED GENERATION

def quick_generate_enhanced(assignment_text):
    """Quick enhanced generation for any assignment"""
    json_output = extract_assignment_structure_advanced(assignment_text)
    if json_output:
        structure = save_assignment_structure_advanced(json_output, 999)  # Temp ID
        if structure:
            prompts = create_advanced_subsection_prompts(structure, assignment_text, "", 999)
            if prompts:
                files = generate_all_assignment_subsections_advanced(prompts, 999)
                if files:
                    return create_professional_docx_from_subsections(999, structure)
    return None

def regenerate_subsection_enhanced(assignment_id, subsection_number):
    """Regenerate specific subsection using enhanced method"""
    try:
        prompts_file = os.path.join("generated_assignments", f"assignment_{assignment_id}", "subsection_prompts.json")
        with open(prompts_file, 'r') as f:
            prompts = json.load(f)
        
        if 0 <= subsection_number - 1 < len(prompts):
            subsection_data = prompts[subsection_number - 1]
            filename = generate_single_subsection_advanced(subsection_data, assignment_id)
            if filename:
                # Reload structure and regenerate DOCX
                structure_file = os.path.join("generated_assignments", f"assignment_{assignment_id}", "assignment_structure.json")
                with open(structure_file, 'r') as f:
                    structure_data = json.load(f)
                return create_professional_docx_from_subsections(assignment_id, structure_data)
        return None
    except Exception as e:
        logger.error(f"Error regenerating subsection: {e}")
        return None

# if __name__ == "__main__":
#     import uvicorn
#     uvicorn.run(app, host="0.0.0.0", port=8000)

active_sessions = {}
session_lock = threading.Lock()

class UserSession:
    def __init__(self):
        self.session_id = str(uuid.uuid4())
        self.temp_dir = tempfile.mkdtemp(prefix=f'user_session_{self.session_id}_')
        self.created_at = datetime.now()
    def cleanup(self):
        try:
            shutil.rmtree(self.temp_dir, ignore_errors=True)
        except Exception as e:
            logger.warning(f"Failed to cleanup session {self.session_id}: {e}")

def create_user_session():
    session = UserSession()
    with session_lock:
        active_sessions[session.session_id] = session
    return session

def cleanup_user_session(session_id):
    with session_lock:
        session = active_sessions.pop(session_id, None)
        if session:
            session.cleanup()

@app.get("/health")
async def health_check():
    return {
        "status": "healthy",
        "active_sessions": len(active_sessions),
        "timestamp": datetime.now().isoformat()
    }

